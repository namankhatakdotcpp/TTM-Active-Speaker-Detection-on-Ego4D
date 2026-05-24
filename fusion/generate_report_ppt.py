"""
Report-accurate presentation for Group 42 CS671 end-semester report.
White background · minimal design · 20 slides · ~15 min

Run:
    python generate_report_ppt.py
Output:
    presentation/TTM_Group42_Final.pptx
"""

import io, os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
import numpy as np
from matplotlib.patches import FancyArrowPatch

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── colour palette (white-bg) ──────────────────────────────────────────────────
WHITE  = RGBColor(0xFF,0xFF,0xFF)
BLACK  = RGBColor(0x1A,0x1A,0x1A)
G90    = RGBColor(0xF5,0xF7,0xFA)   # card fill
G70    = RGBColor(0xDC,0xE1,0xE9)   # border / divider
G40    = RGBColor(0x6B,0x78,0x88)   # secondary text
NAVY   = RGBColor(0x0F,0x2B,0x5B)   # dark navy  (header bars)
BLUE   = RGBColor(0x16,0x5D,0xB6)   # accent blue
TEAL   = RGBColor(0x00,0x96,0x88)   # teal
ORANGE = RGBColor(0xF4,0x73,0x10)
GREEN  = RGBColor(0x2E,0x7D,0x32)
RED    = RGBColor(0xC6,0x28,0x28)
GOLD   = RGBColor(0xF5,0xA8,0x00)

# matplotlib hex
MN = "#0F2B5B"; MB = "#165DB6"; MT = "#009688"; MO = "#F47310"
MG = "#2E7D32"; MR = "#C62828"; MW = "#FFFFFF"; MLG= "#F5F7FA"
MDG= "#6B7888"; MBB= "#EBF2FC"; MGB= "#E8F5E9"; MOB= "#FFF3E0"

SW, SH = 13.33, 7.5
HERE   = os.path.dirname(os.path.abspath(__file__))
OUT    = os.path.join(HERE, "presentation")
os.makedirs(OUT, exist_ok=True)


# ── pptx helpers ──────────────────────────────────────────────────────────────
def _rgb(c):
    if isinstance(c, RGBColor): return c
    if isinstance(c,str) and c.startswith("#"):
        h=c.lstrip("#")
        return RGBColor(int(h[:2],16),int(h[2:4],16),int(h[4:6],16))
    return c

def bg(slide, c=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(c)

def rct(slide,l,t,w,h,fill=WHITE,line=None,lw=Pt(1)):
    sh=slide.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=_rgb(fill)
    if line: sh.line.color.rgb=_rgb(line); sh.line.width=lw
    else:    sh.line.fill.background()
    return sh

def txb(slide,text,l,t,w,h,sz=11,bold=False,italic=False,
        col=BLACK,align=PP_ALIGN.LEFT,wrap=True):
    box=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=box.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=_rgb(col)
    return box

def add_p(tf,text,sz=10,bold=False,italic=False,col=BLACK,align=PP_ALIGN.LEFT):
    p=tf.add_paragraph(); p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=_rgb(col)
    return p

def notes(slide,text):
    slide.notes_slide.notes_text_frame.text=text

def f2b(fig):
    buf=io.BytesIO()
    fig.savefig(buf,format="png",dpi=160,bbox_inches="tight",facecolor=fig.get_facecolor())
    buf.seek(0); plt.close(fig); return buf

# ── slide chrome ─────────────────────────────────────────────────────────────
FOOTER_TXT = "CS671 — Deep Learning and Its Applications  ·  Group 42  ·  IIT Kanpur  ·  May 2026"

def hdr(slide,title,n,total=20):
    rct(slide,0,0,SW,0.07,fill=NAVY)
    txb(slide,title,0.45,0.1,11.2,0.6,sz=22,bold=True,col=WHITE)
    rct(slide,12.3,0.1,0.65,0.38,fill=G90,line=G70)
    txb(slide,f"{n}/{total}",12.3,0.12,0.65,0.35,sz=9,col=G40,align=PP_ALIGN.CENTER)
    rct(slide,0.45,0.78,SW-0.9,0.025,fill=G70)

def ftr(slide):
    rct(slide,0,7.32,SW,0.025,fill=G70)
    txb(slide,FOOTER_TXT,0.45,7.34,SW-0.9,0.2,sz=7.5,col=G40,align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  CHARTS
# ══════════════════════════════════════════════════════════════════════════════

def _ax(ax,xlabel="",ylabel="",title="",grid=True):
    ax.set_facecolor(MW)
    for sp in ax.spines.values(): sp.set_color("#DCE1E9")
    ax.tick_params(colors=MDG,labelsize=8)
    if xlabel: ax.set_xlabel(xlabel,fontsize=9,color=MDG)
    if ylabel: ax.set_ylabel(ylabel,fontsize=9,color=MDG)
    if title:  ax.set_title(title,fontsize=10,fontweight="bold",color=MN,pad=6)
    if grid:   ax.grid(True,alpha=0.35,color="#DCE1E9",linewidth=0.7)


# ── Fig 1: What we do vs others ───────────────────────────────────────────────
def fig_comparison_approach():
    """Side-by-side approach comparison diagram."""
    fig = plt.figure(figsize=(13,5.5),facecolor=MW)
    ax  = fig.add_axes([0,0,1,1])
    ax.set_facecolor(MW); ax.axis("off")
    ax.set_xlim(0,13); ax.set_ylim(0,5.5)

    # title
    ax.text(6.5,5.25,"How TTMFusionModel Differs from Prior Work",
            ha="center",va="center",fontsize=13,fontweight="bold",color=MN)

    # ── prior approaches (left column) ───────────────────────────────────────
    def prior_box(cx,cy,fc,ec,title,desc):
        r=mpatches.FancyBboxPatch((cx-1.7,cy-0.55),3.4,1.1,
            boxstyle="round,pad=0.08",lw=1.4,edgecolor=ec,facecolor=fc)
        ax.add_patch(r)
        ax.text(cx,cy+0.22,title,ha="center",va="center",
                fontsize=9,fontweight="bold",color=MN)
        ax.text(cx,cy-0.17,desc,ha="center",va="center",
                fontsize=7.5,color=MDG,style="italic")

    prior = [
        (2.2,4.4, "#FBE9E7","#E64A19","Audio-only",        "MFCC/log-mel + LSTM/SVM\nIgnores visual cues"),
        (2.2,3.0, "#E3F2FD","#1565C0","Video-only",         "FaceNet/CNN classification\nIgnores prosody & tone"),
        (2.2,1.6, "#F3E5F5","#6A1B9A","Late Fusion",         "Separate models → score combine\nNo cross-modal interaction"),
        (2.2,0.4, "#FFF8E1","#F57F17","Ego4D Baseline",      "SlowFast + audio concat\nLimited context, F1≈0.44"),
    ]
    for args in prior: prior_box(*args)

    # column label
    ax.text(2.2,5.0,"Prior / Baseline Approaches",ha="center",va="center",
            fontsize=10,fontweight="bold",color=MR)
    ax.text(2.2,4.78,"(Limitations)",ha="center",va="center",
            fontsize=8.5,color=MR,style="italic")

    # ── VS divider ────────────────────────────────────────────────────────────
    ax.axvline(5.35,ymin=0.02,ymax=0.88,color=MDG,lw=1.5,ls="--")
    ax.text(5.35,2.6,"VS",ha="center",va="center",fontsize=14,
            fontweight="bold",color=MDG,
            bbox=dict(boxstyle="round,pad=0.3",facecolor=MW,edgecolor=MDG,lw=1.5))

    # ── our approach (right column) ───────────────────────────────────────────
    ours = [
        (9.5,4.45,"Bidirectional Cross-Modal Attention",
         "Video → Audio KV  AND  Audio → Video KV\nShared weights → zero extra parameters",
         MBB,MB),
        (9.5,3.1,"Learned No-Audio Token",
         "512-d trainable vector replaces missing audio\nPrevents degenerate uniform attention scores",
         MGB,MG),
        (9.5,1.75,"2-Layer BiLSTM + LayerNorm",
         "Captures temporal context over T=16 clip windows\n512-d output with stable hidden states",
         "#EDE7F6","#6A1B9A"),
        (9.5,0.4,"Bahdanau Additive Attention",
         "Soft-weights over T time steps → 512-d context vector\nInterpretable: shows which clips drove prediction",
         MOB,MO),
    ]
    for (cx,cy,label,desc,fc,ec) in ours:
        r=mpatches.FancyBboxPatch((cx-2.85,cy-0.55),5.7,1.1,
            boxstyle="round,pad=0.08",lw=1.8,edgecolor=ec,facecolor=fc)
        ax.add_patch(r)
        ax.text(cx,cy+0.22,label,ha="center",va="center",
                fontsize=9.5,fontweight="bold",color=MN)
        ax.text(cx,cy-0.17,desc,ha="center",va="center",
                fontsize=7.5,color=MDG)

    # column label + result pill
    ax.text(9.5,5.0,"TTMFusionModel (Ours)",ha="center",va="center",
            fontsize=10,fontweight="bold",color=MB)
    r2=mpatches.FancyBboxPatch((7.85,4.73),3.3,0.32,
        boxstyle="round,pad=0.05",lw=1.5,edgecolor=MG,facecolor=MGB)
    ax.add_patch(r2)
    ax.text(9.5,4.90,"AUC-ROC=0.847  ·  F1=0.765  ·  +73% over baseline",
            ha="center",va="center",fontsize=8.5,fontweight="bold",color=MG)

    # arrows from prior to our innovations
    arrow_pairs = [(3.95,4.4,6.65,4.45),(3.95,3.0,6.65,3.1),
                   (3.95,1.6,6.65,1.75),(3.95,0.4,6.65,0.4)]
    for (x1,y1,x2,y2) in arrow_pairs:
        ax.annotate("",xy=(x2,y2),xytext=(x1,y1),
            arrowprops=dict(arrowstyle="-|>",color=MDG,lw=1.2,
                connectionstyle="arc3,rad=0.0"))

    return f2b(fig)


# ── Fig 2: Model architecture ─────────────────────────────────────────────────
def fig_arch():
    fig = plt.figure(figsize=(13,5.2),facecolor=MW)
    ax  = fig.add_axes([0,0,1,1])
    ax.set_facecolor(MW); ax.axis("off")
    ax.set_xlim(0,13); ax.set_ylim(0,5.2)

    def box(cx,cy,w,h,fc,ec,lbl,sub="",lfs=9,sfs=7.5):
        r=mpatches.FancyBboxPatch((cx-w/2,cy-h/2),w,h,
            boxstyle="round,pad=0.08",lw=1.5,edgecolor=ec,facecolor=fc)
        ax.add_patch(r)
        ax.text(cx,cy+(0.1 if sub else 0),lbl,ha="center",va="center",
                fontsize=lfs,fontweight="bold",color=MN)
        if sub: ax.text(cx,cy-0.22,sub,ha="center",va="center",
                        fontsize=sfs,color=MDG,style="italic")

    def arr(x1,y1,x2,y2,lbl="",ly=0.12):
        ax.annotate("",xy=(x2,y2),xytext=(x1,y1),
            arrowprops=dict(arrowstyle="-|>",color=MDG,lw=1.6))
        if lbl: ax.text((x1+x2)/2,(y1+y2)/2+ly,lbl,ha="center",va="bottom",
                        fontsize=7.5,color=MDG,style="italic")

    # input labels
    ax.text(0.1,3.9,"Video\nClips",ha="center",va="center",fontsize=8,color=MDG)
    ax.text(0.1,2.1,"Audio\n16kHz",ha="center",va="center",fontsize=8,color=MDG)
    arr(0.4,3.9,0.6,3.9); arr(0.4,2.1,0.6,2.1)

    # encoders
    box(1.3,3.9,1.3,0.65,MBB,MB,"I3D-R50","Frozen\n512-d")
    box(1.3,2.1,1.3,0.65,MGB,MG,"ResNet-18","Frozen\n512-d")
    arr(1.95,3.9,2.2,3.9,"[B,T,512]")
    arr(1.95,2.1,2.2,2.1,"[B,T,512]")

    # projections
    box(3.0,3.9,1.3,0.6,MBB,MB,"Video Proj","Lin+LN+ReLU\n→[B,T,256]")
    box(3.0,2.1,1.3,0.6,MGB,MG,"Audio Proj","Lin+LN+ReLU\n→[B,T,256]")
    arr(3.65,3.9,3.9,3.9)
    arr(3.65,2.1,3.9,2.1)

    # cross-modal attention (tall box spanning both streams)
    r=mpatches.FancyBboxPatch((4.0,1.45),2.6,3.1,
        boxstyle="round,pad=0.1",lw=2.0,edgecolor=MO,facecolor=MOB)
    ax.add_patch(r)
    ax.text(5.3,3.55,"Bidirectional",ha="center",va="center",
            fontsize=9.5,fontweight="bold",color=MN)
    ax.text(5.3,3.25,"Cross-Modal",ha="center",va="center",
            fontsize=9.5,fontweight="bold",color=MN)
    ax.text(5.3,2.95,"Attention",ha="center",va="center",
            fontsize=9.5,fontweight="bold",color=MN)
    ax.text(5.3,2.58,"L=2, 4 heads",ha="center",va="center",
            fontsize=8,color=MDG,style="italic")
    ax.text(5.3,2.28,"shared weights",ha="center",va="center",
            fontsize=8,color=MDG,style="italic")
    ax.text(5.3,1.98,"→ [B,T,512]",ha="center",va="center",
            fontsize=8,color=MO,fontweight="bold")

    # curved arrows into cross-attn
    ax.annotate("",xy=(4.07,3.55),xytext=(3.65,3.9),
        arrowprops=dict(arrowstyle="-|>",color=MDG,lw=1.4,
            connectionstyle="arc3,rad=-0.3"))
    ax.annotate("",xy=(4.07,2.45),xytext=(3.65,2.1),
        arrowprops=dict(arrowstyle="-|>",color=MDG,lw=1.4,
            connectionstyle="arc3,rad=0.3"))

    # BiLSTM
    box(8.0,3.0,1.6,0.7,MBB,MB,"BiLSTM","2-layer, h=256\nLayerNorm\n→[B,T,512]",sfs=7)
    arr(6.6,3.0,7.2,3.0,"[B,T,512]")

    # Bahdanau
    box(10.0,3.0,1.6,0.7,"#EDE7F6","#6A1B9A","Bahdanau\nAttn","Additive soft\n→[B,512]")
    arr(8.8,3.0,9.2,3.0,"[B,T,512]")

    # Classifier
    box(11.9,3.0,1.4,0.6,"#FCE4EC",MR,"Classifier","Dropout(0.5)\nLinear(512,2)")
    arr(10.8,3.0,11.2,3.0,"[B,512]")

    # output
    ax.text(12.88,3.0,"TTM\nlogits",ha="center",va="center",
            fontsize=9,fontweight="bold",color=MB)
    arr(12.6,3.0,12.73,3.0)

    # no-audio token note
    r=mpatches.FancyBboxPatch((4.1,0.1),2.5,0.65,
        boxstyle="round,pad=0.06",lw=1.2,edgecolor=MDG,facecolor=MLG)
    ax.add_patch(r)
    ax.text(5.35,0.49,"★  Learned no-audio token (512-d)",ha="center",va="center",
            fontsize=8.5,fontweight="bold",color=MN)
    ax.text(5.35,0.25,"replaces all-zero audio windows",ha="center",va="center",
            fontsize=7.5,color=MDG,style="italic")
    ax.annotate("",xy=(5.35,1.45),xytext=(5.35,0.75),
        arrowprops=dict(arrowstyle="-|>",color=MDG,lw=1.2,ls="dashed"))

    ax.text(6.5,4.95,"TTMFusionModel — Architecture Overview",
            ha="center",va="center",fontsize=12,fontweight="bold",color=MN)
    ax.text(6.5,4.68,"4,741,155 total parameters  ·  I3D-R50 + ResNet-18 encoders kept frozen",
            ha="center",va="center",fontsize=8.5,color=MDG)

    return f2b(fig)


# ── Fig 3: Dataset distribution ───────────────────────────────────────────────
def fig_dataset():
    fig,axes=plt.subplots(1,2,figsize=(12,4.0),facecolor=MW)
    fig.subplots_adjust(wspace=0.35,left=0.06,right=0.97,top=0.85,bottom=0.12)

    # Class distribution
    ax=axes[0]
    splits=["Full Ego4D\n636k clips","Train split\n55k clips","Val split\n15k clips","Unseen test\n636k clips"]
    ttm   =[5.53,  50, 50,  5.53]
    non   =[94.47, 50, 50, 94.47]
    x=np.arange(len(splits))
    b1=ax.bar(x,non,color=MB,alpha=0.8,edgecolor="none",label="non-TTM")
    b2=ax.bar(x,ttm,bottom=non,color=MO,alpha=0.9,edgecolor="none",label="TTM")
    for xi,tv in zip(x,ttm):
        ax.text(xi,non[xi]+tv/2,f"{tv:.1f}%",ha="center",va="center",
                fontsize=9,fontweight="bold",color=MW)
    ax.set_xticks(x); ax.set_xticklabels(splits,fontsize=8.5)
    ax.set_ylim(0,120); ax.set_ylabel("Percentage (%)",fontsize=9,color=MDG)
    _ax(ax,title="Class Distribution Across Splits",grid=False)
    ax.legend(fontsize=9,edgecolor="#DCE1E9")

    # Model comparison (from results table in report)
    ax2=axes[1]
    models=["Baseline\n(concat+FocalLoss)","Proposed\n@t=0.50","Proposed\n@t=0.139 (val)","Proposed\n@t=0.20 (unseen)"]
    f1s   =[0.44, 0.628, 0.761, 0.765]
    prec  =[0.850,0.837, 0.639, 0.705]
    rec   =[0.300,0.502, 0.940, 0.837]
    bar_w =0.22
    xp=np.arange(len(models))
    bars_f=ax2.bar(xp-bar_w,f1s, bar_w,label="F1",    color=MB,  edgecolor="none",alpha=0.85)
    bars_p=ax2.bar(xp,       prec,bar_w,label="Precision",color=MG,edgecolor="none",alpha=0.85)
    bars_r=ax2.bar(xp+bar_w, rec, bar_w,label="Recall",color=MO, edgecolor="none",alpha=0.85)
    for bars in [bars_f,bars_p,bars_r]:
        for b in bars:
            ax2.text(b.get_x()+b.get_width()/2,b.get_height()+0.008,
                     f"{b.get_height():.2f}",ha="center",va="bottom",
                     fontsize=7,color=MN,fontweight="bold")
    ax2.set_xticks(xp); ax2.set_xticklabels(models,fontsize=8)
    ax2.set_ylim(0,1.18)
    _ax(ax2,ylabel="Score",title="Model Performance Comparison (from Report)")
    ax2.legend(fontsize=8,edgecolor="#DCE1E9")
    ax2.axhline(0.75,color=MG,ls="--",lw=1.2,alpha=0.6)
    ax2.text(3.2,0.77,"0.75 target",fontsize=7,color=MG,style="italic")

    return f2b(fig)


# ── Fig 4: Training history ────────────────────────────────────────────────────
def fig_training():
    hist=[
        (1, 0.761,0.696,0.684,0.640),
        (2, 0.637,0.642,0.730,0.648),
        (3, 0.580,0.632,0.773,0.715),
        (5, 0.561,0.612,0.807,0.723),
        (7, 0.538,0.635,0.812,0.743),
        (9, 0.521,0.649,0.828,0.628),  # best
        (11,0.482,0.703,0.808,0.736),
        (13,0.455,0.704,0.793,0.752),
        (17,0.380,0.851,0.754,0.697),
        (18,0.367,0.790,0.768,0.720),
    ]
    ep =[h[0] for h in hist]
    trl=[h[1] for h in hist]
    val=[h[2] for h in hist]
    ap =[h[3] for h in hist]
    f1 =[h[4] for h in hist]
    lr_changes=[(13,"×0.5 LR"),(17,"×0.5 LR")]

    fig,axes=plt.subplots(1,3,figsize=(13,3.8),facecolor=MW)
    fig.subplots_adjust(wspace=0.35,left=0.05,right=0.98,top=0.85,bottom=0.16)

    # loss
    axes[0].plot(ep,trl,"-o",color=MB,lw=1.8,ms=3.5,label="Train Loss")
    axes[0].plot(ep,val,"-s",color=MO,lw=1.8,ms=3.5,label="Val Loss")
    axes[0].axvline(9,color=MG,ls="--",lw=1.5,alpha=0.8,label="Best (ep.9)")
    axes[0].fill_between(ep,trl,val,alpha=0.06,color=MB)
    _ax(axes[0],"Epoch","Loss","Loss vs Epoch")
    axes[0].legend(fontsize=8,edgecolor="#DCE1E9")

    # val AP and F1
    axes[1].plot(ep,ap,"-o",color=MG,  lw=1.8,ms=3.5,label="Val AP")
    axes[1].plot(ep,f1,"-s",color=MO, lw=1.8,ms=3.5,label="Val F1 @0.5")
    axes[1].axvline(9,color=MG,ls="--",lw=1.5,alpha=0.8,label="Best AP=0.828")
    axes[1].axhline(0.828,color=MG,ls=":",lw=1,alpha=0.5)
    axes[1].text(9.3,0.835,"AP=0.828★",fontsize=7.5,color=MG,fontweight="bold")
    _ax(axes[1],"Epoch","Score","Val AP & F1")
    axes[1].legend(fontsize=8,edgecolor="#DCE1E9")

    # LR schedule
    lr_sched=[(1,3e-4),(12,3e-4),(13,1.5e-4),(16,1.5e-4),(17,7.5e-5),(18,7.5e-5)]
    lx=[s[0] for s in lr_sched]; ly=[s[1] for s in lr_sched]
    axes[2].step(lx,ly,where="post",color=MB,lw=2.0)
    axes[2].fill_between(lx,ly,step="post",alpha=0.12,color=MB)
    for (ep_c,lbl) in lr_changes:
        axes[2].axvline(ep_c,color=MO,ls="--",lw=1.2,alpha=0.8)
        axes[2].text(ep_c+0.2,2.0e-4,lbl,fontsize=7.5,color=MO,rotation=90)
    axes[2].set_yscale("log")
    _ax(axes[2],"Epoch","LR","LR Schedule (ReduceLROnPlateau)")

    return f2b(fig)


# ── Fig 5: Confusion matrices (from report, optimal thresholds) ───────────────
def fig_confusion():
    fig,axes=plt.subplots(1,2,figsize=(12,4.8),facecolor=MW)
    fig.subplots_adjust(wspace=0.4,left=0.05,right=0.97,top=0.85,bottom=0.08)

    cms=[
        (np.array([[172,267],[30,472]]),
         "Val Split — balanced (941 win.)\nThreshold = 0.139  →  F1=0.761  Recall=94.0%"),
        (np.array([[4147,2243],[1041,5349]]),
         "Unseen Tracks — balanced (12,780 win.)\nThreshold = 0.20  →  F1=0.765  Recall=83.7%"),
    ]
    for ax,(cm,title) in zip(axes,cms):
        im=ax.imshow(cm,cmap=plt.cm.Blues,vmin=0,vmax=cm.max())
        plt.colorbar(im,ax=ax,fraction=0.04,pad=0.03)
        ticks=["non-TTM","TTM"]
        ax.set_xticks([0,1]); ax.set_yticks([0,1])
        ax.set_xticklabels([f"Pred\n{t}" for t in ticks],fontsize=10,color=MN)
        ax.set_yticklabels([f"Actual\n{t}" for t in ticks],fontsize=10,color=MN)
        for (i,j),v in np.ndenumerate(cm):
            pct=100*v/cm[i].sum()
            ax.text(j,i,f"{v:,}\n({pct:.1f}%)",ha="center",va="center",
                    fontsize=11 if cm.max()<600 else 9,fontweight="bold",
                    color=MW if v>cm.max()*0.5 else MN)
        ax.set_title(title,fontsize=9,color=MN,pad=8)
        for sp in ax.spines.values(): sp.set_color("#DCE1E9")
        ax.tick_params(colors=MDG)

    return f2b(fig)


# ── Fig 6: Parameter table chart ──────────────────────────────────────────────
def fig_params():
    comps=["Video Proj\n(Lin+LN+ReLU)","Audio Proj\n(Lin+LN+ReLU)",
           "Cross-Modal Attn\n(L=2, 4 heads)","BiLSTM\n(2-layer)+LN",
           "Bahdanau Attn","Classifier\n(Drop+Linear)","No-audio\ntoken"]
    params=[131840,131840,1578496,2102272,263169,1026,512]
    colors=[MB,MT,"#F47310","#6A1B9A",MG,MR,MDG]

    fig,ax=plt.subplots(figsize=(10,3.5),facecolor=MW)
    y=np.arange(len(comps))
    bars=ax.barh(y,params,color=colors,edgecolor="none",height=0.6,alpha=0.88)
    for bar,v in zip(bars,params):
        ax.text(bar.get_width()+12000,bar.get_y()+bar.get_height()/2,
                f"{v:,}",va="center",fontsize=9,color=MN,fontweight="bold")
    ax.set_yticks(y); ax.set_yticklabels(comps,fontsize=9)
    ax.set_xlabel("Parameters",fontsize=9,color=MDG)
    _ax(ax,title="Parameter Count per Component  (Total: 4,741,155)")
    ax.tick_params(axis="y",colors=MN)
    ax.set_xlim(0,2.5e6)
    ax.axvline(4741155/len(comps),color=MDG,ls=":",lw=1,alpha=0.4)
    # total annotation
    ax.text(2.2e6,3.2,"Total: 4.74M",fontsize=11,fontweight="bold",color=MN)
    plt.tight_layout()
    return f2b(fig)


# ── Fig 7: full results comparison (all variants from table) ──────────────────
def fig_results_table():
    """Horizontal table-style bar chart for all 6 model variants from report."""
    variants=[
        "Baseline — concat, FocalLoss,\nno augmentation",
        "Proposed — cross_attn L=2,\nMixUp, CELoss+LS  [t=0.50]",
        "Proposed — val, optimal\nthreshold  [t=0.139]",
        "Proposed — unseen, natural\ndistribution  [t=0.50]",
        "Proposed — unseen, balanced\neval  [t=0.50]",
        "★ Proposed — unseen, balanced,\noptimal threshold  [t=0.20]",
    ]
    f1s  =[0.440,0.628,0.761,0.797,0.668,0.765]
    prec =[0.850,0.837,0.639,0.693,0.857,0.705]
    rec  =[0.300,0.502,0.940,0.937,0.547,0.837]
    apuc =[0.70, 0.828, None, 0.995,0.847, None]

    fig,ax=plt.subplots(figsize=(13,5.0),facecolor=MW)
    y=np.arange(len(variants))
    w=0.25
    bars_f=ax.barh(y+w,   f1s, w, color=MB,  alpha=0.85, edgecolor="none",label="F1")
    bars_p=ax.barh(y,     prec,w, color=MG,  alpha=0.85, edgecolor="none",label="Precision")
    bars_r=ax.barh(y-w,   rec, w, color=MO,  alpha=0.85, edgecolor="none",label="Recall")

    for bars in [bars_f,bars_p,bars_r]:
        for b in bars:
            ax.text(b.get_width()+0.005,b.get_y()+b.get_height()/2,
                    f"{b.get_width():.3f}",va="center",fontsize=8,color=MN)

    # AP/AUC annotations
    for yi,v in zip(y,apuc):
        if v is not None:
            ax.text(1.01,yi+w/2,f"AP/AUC={v:.3f}",va="center",
                    fontsize=7.5,color="#6A1B9A",fontweight="bold")

    # highlight best row
    ax.axhspan(y[-1]-w-0.06,y[-1]+w+0.22,alpha=0.1,color=MG)

    ax.set_yticks(y); ax.set_yticklabels(variants,fontsize=8.5)
    ax.set_xlim(0,1.25); ax.set_xlabel("Score",fontsize=9,color=MDG)
    _ax(ax,title="Full Results Table — All Model Variants (from Report Table 1)")
    ax.legend(fontsize=9,edgecolor="#DCE1E9",loc="lower right")
    ax.axvline(0.75,color=MG,ls="--",lw=1,alpha=0.5)
    ax.tick_params(axis="y",colors=MN)
    plt.tight_layout()
    return f2b(fig)


# ── Fig 8: Augmentation diagram ───────────────────────────────────────────────
def fig_augmentation():
    fig=plt.figure(figsize=(13,3.8),facecolor=MW)
    ax=fig.add_axes([0,0,1,1])
    ax.set_facecolor(MW); ax.axis("off")
    ax.set_xlim(0,13); ax.set_ylim(0,3.8)

    ax.text(6.5,3.55,"Data Augmentation Strategies (Training Only)",
            ha="center",va="center",fontsize=12,fontweight="bold",color=MN)

    augs=[
        (2.2,"MixUp",MOB,MO,
         "p=0.25, λ~Beta(0.4,0.4)",
         "Linear interpolation of two embedding windows\nx̃ = λ·xᵢ + (1−λ)·xⱼ",
         "Improves generalisation; soft labels\nprevent overconfidence"),
        (6.5,"Temporal Masking","#EDE7F6","#6A1B9A",
         "p=0.30, mask ≤ T/4=4 steps",
         "Zero out contiguous time steps in both\nvideo and audio embeddings",
         "Simulates missing clips;\nforces model to be robust to gaps"),
        (10.8,"Feature Noise",MGB,MG,
         "p=0.20, σ=0.01",
         "Gaussian noise N(0, 0.01) added\nto video and audio embeddings",
         "Prevents over-reliance on exact\nembedding values; regularisation"),
    ]
    for (cx,title,fc,ec,param,desc,impact) in augs:
        r=mpatches.FancyBboxPatch((cx-2.0,0.2),4.0,3.0,
            boxstyle="round,pad=0.1",lw=1.8,edgecolor=ec,facecolor=fc)
        ax.add_patch(r)
        rh=mpatches.FancyBboxPatch((cx-2.0,2.82),4.0,0.38,
            boxstyle="round,pad=0.05",lw=0,edgecolor=ec,facecolor=ec)
        ax.add_patch(rh)
        ax.text(cx,3.05,title,ha="center",va="center",fontsize=11,
                fontweight="bold",color=MW)
        ax.text(cx,2.55,param,ha="center",va="center",fontsize=8.5,
                color=MN,fontweight="bold")
        ax.text(cx,2.0, desc, ha="center",va="center",fontsize=8,color=MDG)
        ax.axhline(1.35,xmin=(cx-2.0)/13,xmax=(cx+2.0)/13,
                   color=ec,lw=0.8,alpha=0.5)
        ax.text(cx,0.9, impact,ha="center",va="center",fontsize=8,
                color=MN,style="italic")
        ax.text(cx,0.55,"→ Impact",ha="center",va="center",
                fontsize=8,color=ec,fontweight="bold")

    return f2b(fig)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def s01_title(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    rct(slide,0,0,SW,0.08,fill=NAVY)

    txb(slide,"CS671 — Deep Learning and Its Applications",
        0.5,0.22,12.3,0.5,sz=12,col=G40,italic=True)
    txb(slide,'"Who Is Talking To Me?" — TTM Challenge',
        0.5,0.8,12.3,0.92,sz=32,bold=True,col=NAVY)
    txb(slide,"Multimodal Talking-To-Me Detection via Video–Audio Fusion\n"
              "with Bidirectional Cross-Modal Attention",
        0.5,1.78,12.3,0.65,sz=15,col=G40)
    rct(slide,0.5,2.58,12.3,0.03,fill=G70)

    # team card
    rct(slide,0.5,2.78,7.8,4.2,fill=G90,line=G70)
    txb(slide,"Team",0.65,2.85,2.0,0.4,sz=10,col=G40,bold=True)
    members=[
        "B22109 Kanika Choudhary","B22031 Aman Sharma",
        "B22249 Vikky Kumar",     "B22281 Vershita Yadav",
        "B22051 Mihir Chandra",   "B22043 Harshit",
        "B22069 Sowmika Rao",     "B23217 Naman",
    ]
    for i,m in enumerate(members):
        col=0 if i<4 else 1
        row=i%4
        x=0.7+col*3.8; y=3.3+row*0.72
        rct(slide,x,y,3.4,0.58,fill=WHITE,line=G70)
        txb(slide,m,x+0.12,y+0.08,3.15,0.42,sz=10.5,col=BLACK)
    txb(slide,"Mentor: Jyoti Nigam  |  IIT Kanpur  |  May 1, 2026",
        0.65,6.68,7.5,0.35,sz=10,col=BLUE)

    # results card
    rct(slide,8.65,2.78,4.2,4.2,fill=NAVY)
    txb(slide,"Key Results",8.85,2.9,3.8,0.45,
        sz=13,bold=True,col=WHITE)
    rct(slide,8.65,3.38,4.2,0.03,fill=BLUE)

    res=[
        ("Val AP",         "0.828",  GOLD),
        ("Val F1",         "0.761",  TEAL),
        ("Unseen F1",      "0.765",  GREEN),
        ("AUC-ROC",        "0.847",  ORANGE),
        ("vs Baseline F1", "+73%",   RED),
    ]
    for i,(k,v,ec) in enumerate(res):
        y=3.52+i*0.72
        rct(slide,8.75,y,4.0,0.6,fill=_rgb("#1A3A6B"))
        txb(slide,k, 8.88,y+0.06,2.0,0.45,sz=10,col=G40)
        txb(slide,v, 11.0,y+0.04,1.55,0.52,sz=18,bold=True,
            col=ec,align=PP_ALIGN.RIGHT)

    ftr(slide)
    return slide


def s02_abstract(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Abstract & Problem Overview",2); ftr(slide)

    rct(slide,0.45,1.0,12.5,2.0,fill=G90,line=G70)
    abstract=(
        "We address the problem of egocentric Talking-To-Me (TTM) detection: "
        "given a first-person video clip and its paired audio, determine whether the visible person "
        "is directly addressing the camera wearer. We propose TTMFusionModel, combining frozen "
        "I3D-R50 video and ResNet-18 audio embeddings through bidirectional cross-modal attention "
        "followed by a 2-layer BiLSTM with Bahdanau attention. To counter the 18:1 class imbalance "
        "in Ego4D, we use a track-level balanced 70k-clip training set, window-level weighted "
        "sampling, and CrossEntropyLoss with label smoothing. The model achieves Val AP = 0.828, "
        "evaluation F1 = 0.797 on the val split and AUC-ROC = 0.847 on the full unseen clips index, "
        "substantially outperforming the naïve concatenation baseline (F1 = 0.44)."
    )
    txb(slide,abstract,0.65,1.08,12.1,1.85,sz=12,col=BLACK)

    # problem facets
    facets=[
        (BLUE,   "Multi-modal complexity",
         "Facial/body orientation alone ambiguous; audio fails in noisy multi-speaker environments"),
        (ORANGE, "Severe class imbalance",
         "Only ~5.5% of Ego4D clips are positive → naïve models predict all-negative"),
        (GREEN,  "Strict evaluation",
         "Must be track-level: same person never in both train and val splits"),
        (RED,    "Temporal sequence",
         "Binary classification over sliding windows of clips encoded from both modalities"),
    ]
    txb(slide,"Problem Challenges",0.45,3.15,12.5,0.4,sz=13,bold=True,col=NAVY)
    for i,(ec,title,desc) in enumerate(facets):
        row,col=divmod(i,2)
        x=0.45+col*6.5; y=3.62+row*1.35
        rct(slide,x,y,6.1,1.18,fill=G90)
        rct(slide,x,y,0.06,1.18,fill=ec)
        txb(slide,title,x+0.2,y+0.08,5.7,0.4,sz=11.5,bold=True,col=ec)
        txb(slide,desc, x+0.2,y+0.5, 5.7,0.6, sz=10.5,col=G40)

    notes(slide,"Read the abstract aloud, then walk through the four challenges. Duration: ~1 min.")
    return slide


def s03_comparison(prs,img):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"What We Do vs Prior Work",3); ftr(slide)
    slide.shapes.add_picture(img,Inches(0.2),Inches(0.95),Inches(12.95),Inches(6.1))
    notes(slide,
        "Left column: prior approaches and their limitations. "
        "Right column: our four key innovations. "
        "The VS divider highlights the shift from single-modality / late-fusion to "
        "joint bidirectional cross-modal reasoning. Duration: ~1.5 min.")
    return slide


def s04_dataset(prs,img):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Dataset — Ego4D TTM Benchmark",4); ftr(slide)

    # stats table
    stats=[
        ("Total annotated clips","636,406"),
        ("Unique person tracks","1,321"),
        ("TTM clips (positive)","35,189  (5.53%)"),
        ("Non-TTM clips (negative)","601,217  (94.47%)"),
        ("Positive : Negative ratio","1 : 17.1"),
        ("Avg. clips per track","~481"),
    ]
    txb(slide,"Dataset Statistics",0.45,1.0,5.5,0.4,sz=12,bold=True,col=NAVY)
    for i,(k,v) in enumerate(stats):
        y=1.45+i*0.58
        rct(slide,0.45,y,5.5,0.52,fill=G90 if i%2==0 else WHITE,line=G70)
        txb(slide,k, 0.6, y+0.06,3.0,0.38,sz=10.5,col=G40)
        txb(slide,v, 3.7, y+0.06,2.1,0.38,sz=10.5,bold=True,col=NAVY,
            align=PP_ALIGN.RIGHT)

    # splits table
    txb(slide,"Derived Splits Used",6.3,1.0,6.6,0.4,sz=12,bold=True,col=NAVY)
    hdr_splits=[("CSV / Split","Clips","Tracks","TTM %","Purpose")]
    splits_data=[
        ("balanced_70k — train","55,408","586","50%","Model training"),
        ("balanced_70k — val",  "14,592","146","50%","Early stopping / HP tuning"),
        ("clips_index.csv (unseen)","636,406","1,321","5.53%","Final unseen evaluation"),
    ]
    for i,(row) in enumerate(hdr_splits+splits_data):
        y=1.45+i*0.6
        fill=NAVY if i==0 else (G90 if i%2==1 else WHITE)
        rct(slide,6.3,y,6.6,0.55,fill=fill,line=G70)
        cols_w=[2.1,0.9,0.8,0.7,2.0]
        cols_x=[6.4,8.55,9.5,10.35,11.1]
        for j,(txt,cw) in enumerate(zip(row,cols_w)):
            tc=WHITE if i==0 else (NAVY if j==0 else G40)
            txb(slide,txt,cols_x[j],y+0.08,cw,0.38,
                sz=9 if i>0 else 9.5,bold=(i==0),col=tc)

    txb(slide,
        "★  Track-level split: every clip from a (video_uid, person_id) pair appears in exactly "
        "one split — speaker-leakage is zero.",
        0.45,5.1,12.5,0.55,sz=11,italic=True,col=BLUE)

    slide.shapes.add_picture(img,Inches(0.2),Inches(5.7),Inches(12.95),Inches(1.45))

    notes(slide,
        "Emphasise the 18:1 imbalance in the raw dataset and how the 70k balanced subset "
        "addresses it. Track-level split ensures no person is in both train and test. Duration: ~1 min.")
    return slide


def s05_architecture(prs,img):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Model Architecture — TTMFusionModel",5); ftr(slide)
    slide.shapes.add_picture(img,Inches(0.2),Inches(0.95),Inches(12.95),Inches(5.6))
    notes(slide,
        "Walk L→R through the 5 learned stages: (1) modality projection, "
        "(2) bidirectional cross-modal attention, (3) BiLSTM+LN, "
        "(4) Bahdanau attention, (5) classifier. "
        "Highlight the no-audio token below. Duration: ~2 min.")
    return slide


def s06_params(prs,img):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Architecture Components & Parameter Count",6); ftr(slide)

    slide.shapes.add_picture(img,Inches(0.25),Inches(1.0),Inches(9.5),Inches(4.2))

    # table on right
    rows=[
        ("Video Proj (Lin+LN+ReLU)","[B,T,256]","131,840"),
        ("Audio Proj (Lin+LN+ReLU)","[B,T,256]","131,840"),
        ("Cross-Modal Attn (L=2,4h)","[B,T,512]","1,578,496"),
        ("BiLSTM (2-layer) + LN",   "[B,T,512]","2,102,272"),
        ("Bahdanau Attention",       "[B,512]",  "263,169"),
        ("Classifier (Drop+Linear)", "[B,2]",    "1,026"),
        ("No-audio token",           "512-d param","512"),
    ]
    txb(slide,"Component",10.0,1.0,2.5,0.35,sz=9.5,bold=True,col=NAVY)
    txb(slide,"Output",12.55,1.0,1.0,0.35,sz=9.5,bold=True,col=NAVY,align=PP_ALIGN.RIGHT)

    for i,(comp,shape,param) in enumerate(rows):
        y=1.4+i*0.6
        rct(slide,10.0,y,3.25,0.55,fill=G90 if i%2==0 else WHITE,line=G70)
        txb(slide,comp,  10.1,y+0.06,1.75,0.38,sz=9,col=BLACK)
        txb(slide,shape, 11.88,y+0.06,1.2,0.38,sz=8.5,col=G40,align=PP_ALIGN.CENTER)
        txb(slide,param, 12.55,y+0.06,0.65,0.38,sz=9,col=BLUE,bold=True,align=PP_ALIGN.RIGHT)

    # total row
    rct(slide,10.0,1.4+7*0.6,3.25,0.52,fill=NAVY)
    txb(slide,"Total",10.1,1.4+7*0.6+0.06,1.6,0.38,sz=11,bold=True,col=WHITE)
    txb(slide,"4,741,155",11.6,1.4+7*0.6+0.04,1.6,0.42,sz=13,bold=True,
        col=GOLD,align=PP_ALIGN.RIGHT)

    notes(slide,
        "Cross-modal attention is the largest single component (1.58M params). "
        "BiLSTM adds 2.1M. Total is 4.74M — compact enough for GPU inference. "
        "Duration: ~1 min.")
    return slide


def s07_training(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Training Configuration",7); ftr(slide)

    config=[
        ("Optimiser",       "AdamW","Decoupled weight decay; better generalisation than Adam"),
        ("Learning rate",   "3×10⁻⁴ → ReduceLROnPlateau\n(×0.5, patience=3)",
         "Halved at epochs 13 & 17 when Val AP stagnated"),
        ("Weight decay",    "5×10⁻⁴","Stronger L2 regularisation on the 70k balanced set"),
        ("Dropout",         "0.5","BiLSTM inter-layer + before classifier"),
        ("Batch size",      "256","Stabilises WeightedRandomSampler class-ratio estimates"),
        ("Loss function",   "CrossEntropyLoss\nlabel_smoothing=0.1",
         "Prevents overconfidence; no double-correction with sampler active"),
        ("Class balancing", "WeightedRandomSampler\n(50:50 window level)",
         "Compensates 17:1 imbalance throughout all epochs"),
        ("Early stopping",  "Val AP, patience=10","Threshold-agnostic metric; unbiased for imbalanced classes"),
        ("Window / Stride", "T=16; train=4, val=16, eval=1",
         "Dense training overlap; unbiased val; maximum eval support"),
        ("Best epoch",      "9 / 18  (Val AP=0.828)","Early stopping triggered; 10 more epochs of no improvement"),
        ("Hardware",        "PyTorch · NVIDIA GPU (CUDA 7)","~15 min/epoch on 70k balanced clips"),
        ("Fusion type",     "cross_attn, bidir., L=2","+3% Val AP over concatenation baseline"),
    ]

    for i,(k,v,rat) in enumerate(config):
        row,col=divmod(i,2)
        x=0.45+col*6.5; y=1.05+row*1.05
        rct(slide,x,y,6.1,0.92,fill=G90)
        rct(slide,x,y,2.0,0.92,fill=_rgb("#EBF2FC"))
        txb(slide,k,x+0.1,y+0.05,1.8,0.38,sz=10,bold=True,col=NAVY)
        txb(slide,v,x+2.1,y+0.04,1.85,0.45,sz=9.5,bold=True,col=BLUE)
        txb(slide,rat,x+4.0,y+0.05,2.0,0.82,sz=8.5,col=G40,italic=True)

    notes(slide,
        "Walk through each hyperparameter and its rationale — especially why WeightedRandomSampler "
        "was chosen over FocalLoss (already balanced batches, FocalLoss inflates FP). Duration: ~1.5 min.")
    return slide


def s08_methodology(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Methodology — End-to-End Pipeline",8); ftr(slide)

    # pipeline steps (from report §2.1)
    steps=[
        (MB,  "①\nData\nIngestion",
         "clips_index.csv indexed\n70k balanced subset: 50:50 TTM/non-TTM\nTrack-level 80/20 split"),
        (MT,  "②\nEmbedding\nExtraction",
         "I3D-R50 → 512-d video .npy\nResNet-18 on 128-band log-mel\n→ 512-d audio .npy (offline)"),
        ("#6A1B9A","③\nWindowed\nDataset",
         "T=16 sliding windows; majority-vote label\nTrain stride=4 (75% overlap, 4× windows)\nEval stride=1 (dense, 14,234 win.)"),
        (MO,  "④\nTraining",
         "AdamW + WeightedRandomSampler\nCrossEntropyLoss (smoothing=0.1)\nMixUp + masking + noise"),
        (MR,  "⑤\nEvaluation",
         "Best ckpt (ep.9, AP=0.828)\nArchitecture auto-detected\nBalanced val + full unseen test"),
    ]

    for i,(ec,title,desc) in enumerate(steps):
        x=0.55+i*2.55
        r=mpatches.FancyBboxPatch
        rct(slide,x,1.05,2.25,5.1,fill=G90,line=_rgb(ec))
        rct(slide,x,1.05,2.25,0.52,fill=_rgb(ec))
        txb(slide,title.replace("\n"," "),x+0.1,1.1,2.05,0.4,sz=11,bold=True,col=WHITE)
        txb(slide,desc,x+0.1,1.65,2.05,4.4,sz=9.5,col=BLACK)
        if i<4:
            txb(slide,"→",x+2.3,3.35,0.22,0.4,sz=16,bold=True,col=G40,
                align=PP_ALIGN.CENTER)

    # preprocessing detail
    rct(slide,0.45,6.28,12.5,0.65,fill=_rgb(MBB),line=_rgb(MB))
    txb(slide,
        "Preprocessing:  Video: I3D-R50 (frozen) → 512-d/clip  ·  "
        "Audio: 16kHz → 128-band log-mel (FFT=1024, hop=10ms) → ResNet-18 (frozen) → 512-d/clip  ·  "
        "Coverage fixed 6%→100% by repairing clip-ID path mapping",
        0.6,6.32,12.1,0.58,sz=9.5,col=NAVY)

    notes(slide,
        "Five pipeline steps. Emphasis: step 2 (embedding extraction) is done offline once; "
        "step 4 trains from pre-computed .npy files in ~15 min/epoch. "
        "Audio coverage fix (6%→100%) was the single largest F1 improvement (~11 F1 pts). "
        "Duration: ~1.5 min.")
    return slide


def s09_augmentation(prs,img):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Data Augmentation (Training Only)",9); ftr(slide)
    slide.shapes.add_picture(img,Inches(0.2),Inches(0.98),Inches(12.95),Inches(4.2))

    # key equations
    rct(slide,0.45,5.28,12.5,1.8,fill=G90,line=G70)
    txb(slide,"Key Equations",0.65,5.35,12.0,0.38,sz=11,bold=True,col=NAVY)
    eqs=[
        ("MixUp:",      "x̃ = λ·xᵢ + (1−λ)·xⱼ  ,  ỹ = λ·yᵢ + (1−λ)·yⱼ  ,  λ ~ Beta(0.4, 0.4)"),
        ("Temp. mask:", "zero T/4 contiguous clips in both video & audio embeddings  (simulates missing frames)"),
        ("Feature noise:","ε ~ N(0, 0.01)  added to both v̂ and â at each forward pass"),
    ]
    for i,(k,v) in enumerate(eqs):
        y=5.78+i*0.41
        txb(slide,k,0.65,y,1.6,0.38,sz=9.5,bold=True,col=BLUE)
        txb(slide,v,2.35,y,10.5,0.38,sz=9.5,col=BLACK)

    notes(slide,
        "Three augmentation strategies applied only during training. MixUp uses soft labels, "
        "which is important — do NOT use hard labels with MixUp. "
        "Temporal masking forces the model to classify even when some clips are blank. "
        "Duration: ~1 min.")
    return slide


def s10_equations(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Mathematical Model",10); ftr(slide)

    sections=[
        (0.45,1.0,BLUE,
         "① Modality Projection",
         [("v̂  =  ReLU(LayerNorm(W_v · x_video + b_v))   ∈ ℝᴮˣᵀˣ²⁵⁶",
           "Linear 512→256 + LayerNorm + ReLU; identical for audio â"),]),

        (0.45,2.25,MO,
         "② Bidirectional Cross-Modal Attention (MHA)",
         [("Attn(Q,K,V) = softmax(QKᵀ / √d_k) · V",
           "Multi-head attention, d_k=64, 4 heads, 2 layers"),
          ("v' = v̂ + Attn(v̂, â, â)  ;  a' = â + Attn(â, v̂, v̂)  (shared layer weights)",
           "Bidirectional: each modality attends to the other — zero extra parameters"),
          ("x  = concat[v', a']  ∈ ℝᴮˣᵀˣ⁵¹²",
           "Fused representation fed into BiLSTM"),]),

        (0.45,4.1,MG,
         "③ Bahdanau Additive Attention",
         [("e_t = v⊤ · tanh(W_h · h_t)    ∈ ℝ",
           "Energy score — how important is BiLSTM state h_t?"),
          ("α_t = softmax(e_t)  ∈ ℝᵀ  ;  c = Σ_t (α_t · h_t)  ∈ ℝ⁵¹²",
           "Attention weights (interpretable) → context vector for classifier"),]),

        (6.8,1.0,"#6A1B9A",
         "④ CrossEntropyLoss + Label Smoothing",
         [("LS(p, y) = (1−ε)·CE(p, y) + ε/K ,   ε=0.1, K=2",
           "Soft targets: true class gets 0.9 prob mass, prevents overconfidence"),]),

        (6.8,2.4,MR,
         "⑤ WeightedRandomSampler",
         [("w_i  =  1 / count(class_i)  →  E[batch TTM%] = 50%",
           "Inverse-frequency weights ensure balanced gradient signal throughout training"),]),

        (6.8,3.55,MDG,
         "⑥ ReduceLROnPlateau Schedule",
         [("lr_new = lr × 0.5   if   Val_AP stagnates for 3 consecutive epochs",
           "Triggered at epochs 13 (lr: 3e-4→1.5e-4) and 17 (lr: 1.5e-4→7.5e-5)"),]),
    ]

    for (x,y,ec,title,eqs) in sections:
        rct(slide,x,y,6.2,0.35,fill=_rgb(ec))
        txb(slide,title,x+0.12,y+0.04,5.95,0.28,sz=10,bold=True,col=WHITE)
        h=0.42
        for j,(eq,exp) in enumerate(eqs):
            ey=y+0.38+j*0.57
            rct(slide,x,ey,6.2,0.52,fill=G90)
            txb(slide,eq, x+0.12,ey+0.02,5.9,0.28,sz=8.8,col=BLACK)
            txb(slide,f"▸ {exp}",x+0.12,ey+0.28,5.9,0.22,sz=8,col=G40,italic=True)

    notes(slide,
        "Walk through each equation. Emphasis: (2) cross-modal MHA with shared weights — "
        "this is the innovation. (3) Bahdanau gives interpretable time-step importance. "
        "(4) label smoothing prevents overconfidence on the balanced set. Duration: ~2 min.")
    return slide


def s11_results_table(prs,img):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Results — All Model Variants",11); ftr(slide)
    slide.shapes.add_picture(img,Inches(0.2),Inches(0.98),Inches(12.95),Inches(5.9))
    notes(slide,
        "Six model variants from the report table. Emphasise: (1) baseline F1=0.44, "
        "(2) proposed @t=0.5 F1=0.628, AP=0.828, "
        "(3) optimal threshold on val F1=0.761, "
        "(4) unseen tracks, optimal F1=0.765, Acc=74.3%. Duration: ~1.5 min.")
    return slide


def s12_confusion(prs,img):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Confusion Matrices at Optimal Thresholds",12); ftr(slide)
    slide.shapes.add_picture(img,Inches(0.2),Inches(1.0),Inches(12.95),Inches(5.1))

    txb(slide,
        "Val (t=0.139): 472/502 TTM correctly detected — 94.0% recall  ·  "
        "Unseen (t=0.20): 5,349/6,390 TTM correctly detected — 83.7% recall",
        0.45,6.2,12.5,0.55,sz=11,col=BLUE,italic=True,align=PP_ALIGN.CENTER)

    notes(slide,
        "Left matrix: validation set, threshold tuned to 0.139. Very high recall (94%) "
        "at cost of some precision. Right: unseen tracks, threshold 0.20, recall=83.7%, "
        "F1=0.765. Model is conservative at t=0.5 but strong after threshold calibration. "
        "Duration: ~1 min.")
    return slide


def s13_training_curves(prs,img):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Training Curves — 18 Epochs",13); ftr(slide)
    slide.shapes.add_picture(img,Inches(0.2),Inches(1.0),Inches(12.95),Inches(5.2))

    txb(slide,
        "★ Best checkpoint epoch 9 (Val AP=0.828)  ·  "
        "LR halved at epochs 13 & 17  ·  Early stopping triggered at epoch 18",
        0.45,6.3,12.5,0.4,sz=10.5,col=G40,italic=True,align=PP_ALIGN.CENTER)

    notes(slide,
        "Left: loss diverges after ep.9 — early stop was correct. "
        "Centre: Val AP peaks 0.828 at ep.9; F1@0.5 oscillates 0.63–0.75 "
        "as precision/recall balance shifts. Right: ReduceLROnPlateau halves LR "
        "at ep.13 and ep.17. Duration: ~1 min.")
    return slide


def s14_observations(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Qualitative Observations & Analysis",14); ftr(slide)

    obs=[
        (BLUE,  "Best metric is AP=0.828, not F1@0.50",
         "At default t=0.5 the model is conservative (F1=0.628, Precision=0.837, Recall=0.502). "
         "AP is threshold-agnostic and the correct single-number summary. "
         "Lowering threshold to 0.139 boosts recall 50%→94% at acceptable precision cost."),
        (MO,    "Unseen-track generalisation",
         "At optimal threshold (t=0.20) on fully unseen tracks: F1=0.765, "
         "Accuracy=74.3%, AUC-ROC=0.847. Performance degrades gracefully — not catastrophically."),
        (MG,    "Cross-modal attention vs concatenation",
         "Bidirectional cross-modal attention (L=2, shared weights) improved Val AP "
         "by ~3% over simple concatenation with only 1.58M extra parameters — "
         "capturing simultaneous audio-visual co-occurrences (lip timing + audio peaks)."),
        (MR,    "Audio coverage was the largest single gain",
         "Fixing audio extraction from 6%→100% clip coverage added ~11 F1 points. "
         "ResNet-18 log-mel embeddings carry prosody and timing information "
         "that video alone cannot provide."),
        ("#6A1B9A","Threshold calibration needed",
         "Optimal val threshold is 0.139 (≪ 0.5), indicating the model systematically "
         "underestimates positive probability. Temperature scaling or Platt calibration "
         "would bring the operating threshold closer to 0.5 for deployment."),
        (MDG,   "Track-level isolation verified programmatically",
         "Intersection of (video_uid, person_id) pairs between training and unseen test is empty. "
         "The model has never seen any of the 264 test tracks during training."),
    ]

    for i,(ec,title,desc) in enumerate(obs):
        row,col=divmod(i,2)
        x=0.45+col*6.5; y=1.05+row*2.05
        rct(slide,x,y,6.1,1.85,fill=G90)
        rct(slide,x,y,0.06,1.85,fill=_rgb(ec))
        txb(slide,title,x+0.2,y+0.08,5.7,0.42,sz=11,bold=True,col=_rgb(ec))
        txb(slide,desc,  x+0.2,y+0.52,5.7,1.25,sz=10,col=BLACK)

    notes(slide,
        "Walk through the 6 observations. Key talking points: "
        "(1) AP vs F1 difference, (2) threshold calibration issue, "
        "(3) cross-modal attention gain, (4) audio coverage fix impact. "
        "Duration: ~2 min.")
    return slide


def s15_challenges(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Challenges Faced & Solutions",15); ftr(slide)

    rows=[
        (RED,   "Class Imbalance  (18:1 non-TTM:TTM)",
         "Without WeightedRandomSampler, model collapses to all-negative (accuracy≈94.5%, TTM F1=0.0).",
         "WeightedRandomSampler enforces 50:50 batches. CrossEntropyLoss with label smoothing 0.1 "
         "prevents overconfidence. Window-level accuracy grows from 59% → confirms sampler is working."),
        (ORANGE,"Missing Audio Embeddings",
         "6% of clips initially had no audio .npy. Zero vectors corrupt cross-modal attention "
         "(uniform scores), causing the model to ignore audio entirely.",
         "★ Learned no-audio token: trainable 512-d parameter replaces all-zero audio windows. "
         "Audio coverage fixed to 100% by repairing clip-ID path mapping — adding ~11 F1 points."),
        (BLUE,  "Speaker Leakage in Data Splits",
         "Naïve random split lets the same person appear in both train and val. "
         "Model memorises speaker-specific cues; validation metrics are inflated.",
         "Track-level split: entire (video_uid, person_id) goes to exactly one split. "
         "Zero-leakage verified programmatically — no track overlap between train and unseen test."),
        (GREEN, "Overfitting on 70k Balanced Subset",
         "Train accuracy rises to 90% by epoch 18 while Val AP declines from 0.828 to 0.768. "
         "Model memorises the small balanced subset.",
         "Dropout=0.5 + weight_decay=5e-4 + MixUp + temporal masking + early stopping (patience=10). "
         "Full clips_index_resplit.csv (636k clips) expected to close this gap."),
    ]

    for i,(ec,title,prob,sol) in enumerate(rows):
        y=1.05+i*1.55
        rct(slide,0.45,y,12.5,1.38,fill=G90)
        rct(slide,0.45,y,0.06,1.38,fill=_rgb(ec))
        txb(slide,title,  0.58,y+0.06,12.2,0.38,sz=11.5,bold=True,col=_rgb(ec))
        txb(slide,"Problem: "+prob,0.58,y+0.48,5.9,0.82,sz=9.5,col=G40,italic=True)
        rct(slide,6.65,y+0.38,6.2,0.92,fill=_rgb(MGB) if i%2==0 else _rgb(MBB))
        txb(slide,"Solution: "+sol,6.75,y+0.44,6.0,0.82,sz=9.5,col=_rgb(MN))

    notes(slide,
        "The no-audio token is an original engineering contribution. "
        "The audio coverage fix (+11 F1 pts) and the track-level split are the two "
        "most impactful practical decisions. Duration: ~1 min.")
    return slide


def s16_future(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Future Scope",16); ftr(slide)

    col_data=[
        (BLUE, "Near-Term", [
            "Temperature scaling / Platt calibration — bring operating threshold closer to 0.5",
            "Train on full clips_index_resplit.csv (636k clips, 1,321 tracks) to close the 70k overfitting gap",
            "Multi-person scoring: simultaneous TTM scores for all visible people in one forward pass",
            "Real-time inference at 30fps via TorchScript export on edge GPU",
        ]),
        (ORANGE,"Mid-Term",[
            "Add ASR text stream as a third modality (text + video + audio joint fusion)",
            "Gaze estimation features (head pose, eye contact) to improve precision",
            "Submit to official Ego4D TTM benchmark leaderboard for community comparison",
            "Explore transformer-based temporal model replacing BiLSTM for longer contexts",
        ]),
        (GREEN,"Long-Term",[
            "Self-supervised pre-training on unlabelled egocentric video (Ego4D unlabelled)",
            "Cross-cultural / multi-language TTM generalisation across diverse speakers",
            "On-device deployment in smart glasses / hearing aids / social robots",
            "Federated learning for privacy-preserving TTM training on wearable data",
        ]),
    ]

    for ci,(ec,title,items) in enumerate(col_data):
        x=0.45+ci*4.3
        rct(slide,x,1.0,4.0,6.08,fill=G90)
        rct(slide,x,1.0,4.0,0.45,fill=_rgb(ec))
        txb(slide,title,x+0.1,1.05,3.8,0.35,sz=13,bold=True,col=WHITE)
        for j,item in enumerate(items):
            y=1.55+j*1.35
            rct(slide,x+0.12,y,3.76,1.2,fill=WHITE,line=G70)
            txb(slide,f"• {item}",x+0.22,y+0.08,3.55,1.08,sz=10,col=BLACK)

    notes(slide,
        "Most impactful near-term: threshold calibration and training on full dataset. "
        "Most exciting long-term: wearable deployment. "
        "Duration: ~45 sec.")
    return slide


def s17_applications(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Applications & Real-World Impact",17); ftr(slide)

    apps=[
        ("Assistive\nTechnology",    BLUE,
         "Smart hearing aids that alert the wearer when directly addressed, "
         "even in noisy multi-speaker environments."),
        ("Social\nRobotics",         ORANGE,
         "Robots detect which human is talking to them, enabling natural "
         "turn-taking in multi-person interactions."),
        ("Meeting\nAnalytics",       GREEN,
         "Automatic speaker diarisation and active-listener detection "
         "in video conferencing systems."),
        ("AR / VR\nWearables",       RED,
         "Context-aware overlays activate only when the user is being "
         "addressed, reducing cognitive load."),
        ("Retail &\nHospitality AI", TEAL,
         "Detect customer–staff engagement; measure response time "
         "and service quality from egocentric store cameras."),
        ("Dementia\n& Elder Care",   GOLD,
         "Monitor social engagement; alert caregivers when a patient "
         "is addressed but does not respond."),
    ]

    for i,(title,ec,desc) in enumerate(apps):
        row,col=divmod(i,3)
        x=0.45+col*4.3; y=1.1+row*2.95
        rct(slide,x,y,4.0,2.65,fill=G90)
        rct(slide,x,y,0.06,2.65,fill=ec)
        txb(slide,title.replace("\n"," "),x+0.2,y+0.1,3.65,0.55,
            sz=12,bold=True,col=ec)
        txb(slide,desc,x+0.2,y+0.72,3.65,1.8,sz=10.5,col=BLACK)

    notes(slide,"Broad applicability — accessibility, robotics, analytics, healthcare. Duration: ~45 sec.")
    return slide


def s18_conclusion(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"Conclusion",18); ftr(slide)

    txb(slide,
        "We developed TTMFusionModel, a multimodal TTM detection system "
        "combining frozen I3D-R50 (video) and ResNet-18 (audio) encoders through bidirectional "
        "cross-modal attention, 2-layer BiLSTM temporal modelling, and Bahdanau attention pooling "
        "— trained and evaluated on the large-scale Ego4D benchmark.",
        0.45,1.02,12.5,0.85,sz=13,col=BLACK)

    # highlight metric table
    for i,(k,v,ec) in enumerate([
        ("Val AP",            "0.828",BLUE),
        ("Val F1 @t=0.139",  "0.761",TEAL),
        ("Unseen F1 @t=0.20","0.765",GREEN),
        ("AUC-ROC (unseen)", "0.847",ORANGE),
    ]):
        x=0.45+i*3.15
        rct(slide,x,2.12,2.95,0.95,fill=NAVY)
        txb(slide,k,  x+0.12,2.18,2.7,0.32,sz=9,col=G40)
        txb(slide,v,  x+0.12,2.48,2.7,0.52,sz=20,bold=True,col=_rgb(ec),
            align=PP_ALIGN.CENTER)

    achievements=[
        ("Innovation",
         "Bidirectional cross-modal attention with shared weights gives both modalities "
         "cross-modal context with zero extra parameter overhead."),
        ("Data Rigour",
         "Track-level split + no-audio token + 100% audio coverage fix = "
         "fair, leakage-free, fully multimodal evaluation."),
        ("Baseline Improvement",
         "F1 raised from 0.44 (concat baseline) to 0.765 (+73%) on unseen "
         "tracks at optimal threshold — a substantial margin."),
        ("Interpretability",
         "Bahdanau attention weights reveal which temporal windows drove each "
         "prediction — useful for debugging and downstream analysis."),
    ]

    for i,(tag,desc) in enumerate(achievements):
        y=3.3+i*1.0
        rct(slide,0.45,y,2.0,0.82,fill=BLUE)
        txb(slide,tag,0.55,y+0.12,1.8,0.58,sz=11,bold=True,col=WHITE,
            align=PP_ALIGN.CENTER)
        rct(slide,2.52,y,10.4,0.82,fill=G90)
        txb(slide,desc,2.65,y+0.1,10.15,0.65,sz=11,col=BLACK)

    notes(slide,
        "Four key achievements. Emphasise the +73% F1 improvement over baseline "
        "and the interpretability benefit of Bahdanau attention. Duration: ~1 min.")
    return slide


def s19_references(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide,"References",19); ftr(slide)

    refs=[
        "[1]  Grauman, K. et al. (2022). Ego4D: Around the World in 3,000 Hours of Egocentric Video. "
             "IEEE/CVF CVPR. arXiv:2110.07177",
        "[2]  Carreira, J. & Zisserman, A. (2017). Quo Vadis, Action Recognition? "
             "A New Model and the Kinetics Dataset. CVPR. [I3D-R50] arXiv:1705.07750",
        "[3]  He, K., Zhang, X., Ren, S. & Sun, J. (2016). Deep Residual Learning for Image Recognition. "
             "CVPR. [ResNet-18] arXiv:1512.03385",
        "[4]  Bahdanau, D., Cho, K. & Bengio, Y. (2015). Neural Machine Translation by Jointly "
             "Learning to Align and Translate. ICLR. arXiv:1409.0473",
        "[5]  Vaswani, A. et al. (2017). Attention Is All You Need. NeurIPS. arXiv:1706.03762",
        "[6]  Loshchilov, I. & Hutter, F. (2019). Decoupled Weight Decay Regularization. "
             "ICLR. [AdamW] arXiv:1711.05101",
        "[7]  Zhang, H. et al. (2018). mixup: Beyond Empirical Risk Minimization. "
             "ICLR. arXiv:1710.09412",
        "[8]  Lin, T.-Y. et al. (2020). Focal Loss for Dense Object Detection. "
             "IEEE TPAMI. arXiv:1708.02002",
        "[9]  Müller, R., Kornblith, S. & Hinton, G. (2019). When Does Label Smoothing Help? "
             "NeurIPS. arXiv:1906.02629",
        "[10] Hochreiter, S. & Schmidhuber, J. (1997). Long Short-Term Memory. "
             "Neural Computation, 9(8), 1735–1780.",
    ]

    for i,ref in enumerate(refs):
        y=1.05+i*0.62
        rct(slide,0.45,y,12.5,0.56,fill=G90 if i%2==0 else WHITE,line=G70)
        txb(slide,ref,0.6,y+0.06,12.2,0.46,sz=9.5,col=BLACK)

    notes(slide,"Display only — no need to read aloud. Duration: <15 sec.")
    return slide


def s20_thankyou(prs):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    rct(slide,0,0,SW,0.08,fill=NAVY)

    txb(slide,"Thank You",0.5,1.1,12.3,1.3,sz=56,bold=True,col=NAVY,
        align=PP_ALIGN.CENTER)
    rct(slide,1.5,2.6,10.33,0.04,fill=BLUE)

    txb(slide,"Questions & Discussion",0.5,2.75,12.3,0.65,
        sz=20,italic=True,col=G40,align=PP_ALIGN.CENTER)

    txb(slide,
        '"Who Is Talking To Me?" — TTM Challenge\n'
        "CS671 · Group 42 · IIT Kanpur · May 2026",
        0.5,3.55,12.3,0.75,sz=13,col=G40,align=PP_ALIGN.CENTER)

    # final metrics
    for i,(k,v,ec) in enumerate([
        ("Val AP",       "0.828",BLUE),
        ("Val F1",       "0.761",TEAL),
        ("Unseen F1",    "0.765",GREEN),
        ("AUC-ROC",      "0.847",ORANGE),
        ("vs Baseline",  "+73%", RED),
    ]):
        x=0.95+i*2.3
        rct(slide,x,4.6,2.1,1.05,fill=G90,line=G70)
        rct(slide,x,4.6,2.1,0.3,fill=_rgb(ec))
        txb(slide,k,x+0.1,4.63,1.9,0.23,sz=8.5,col=WHITE)
        txb(slide,v,x+0.1,4.9, 1.9,0.68,sz=19,bold=True,col=_rgb(ec),
            align=PP_ALIGN.CENTER)

    txb(slide,
        "Mentor: Jyoti Nigam  ·  Members: Kanika · Aman · Vikky · Vershita · "
        "Mihir · Harshit · Sowmika · Naman",
        0.5,5.9,12.3,0.45,sz=10,col=G40,align=PP_ALIGN.CENTER)
    txb(slide,"Contact: gehrishiksha@gmail.com",
        0.5,6.4,12.3,0.35,sz=10,col=BLUE,align=PP_ALIGN.CENTER)

    ftr(slide)
    notes(slide,
        "Open for questions. Likely Q: why not use a full transformer? "
        "Answer: cross-modal attention on top of BiLSTM is more parameter-efficient and "
        "less prone to overfit on 70k clips. Duration: remaining time.")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    print("Generating charts from report data…")
    img_cmp  = fig_comparison_approach()   # what we do vs others
    img_arch = fig_arch()
    img_ds   = fig_dataset()
    img_tr   = fig_training()
    img_cm   = fig_confusion()
    img_par  = fig_params()
    img_res  = fig_results_table()
    img_aug  = fig_augmentation()

    print("Building slides…")
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)

    s01_title(prs)                        #  1
    s02_abstract(prs)                     #  2
    s03_comparison(prs, img_cmp)          #  3
    s04_dataset(prs, img_ds)              #  4
    s05_architecture(prs, img_arch)       #  5
    s06_params(prs, img_par)              #  6
    s07_training(prs)                     #  7
    s08_methodology(prs)                  #  8
    s09_augmentation(prs, img_aug)        #  9
    s10_equations(prs)                    # 10
    s11_results_table(prs, img_res)       # 11
    s12_confusion(prs, img_cm)            # 12
    s13_training_curves(prs, img_tr)      # 13
    s14_observations(prs)                 # 14
    s15_challenges(prs)                   # 15
    s16_future(prs)                       # 16
    s17_applications(prs)                 # 17
    s18_conclusion(prs)                   # 18
    s19_references(prs)                   # 19
    s20_thankyou(prs)                     # 20

    out = os.path.join(OUT, "TTM_Group42_Final.pptx")
    prs.save(out)
    print(f"\nSaved ({len(prs.slides)} slides, {os.path.getsize(out)//1024} KB): {out}")

    import subprocess
    try:
        r=subprocess.run(
            ["libreoffice","--headless","--convert-to","pdf","--outdir",OUT,out],
            capture_output=True,text=True,timeout=120)
        if r.returncode==0:
            print(f"PDF: {out.replace('.pptx','.pdf')}")
        else:
            print("LibreOffice not available — open PPTX in PowerPoint/Slides for PDF.")
    except (FileNotFoundError,subprocess.TimeoutExpired):
        print("LibreOffice not found — PPTX saved. Open in PowerPoint → Save As PDF.")


if __name__=="__main__":
    main()
