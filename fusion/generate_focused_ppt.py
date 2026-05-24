"""
Focused 8-slide presentation — Group 42 CS671
Slides: Title · Problem · Dataset Stats · Approach · Prior Work ·
        Challenges · Results · Future Work

Output: presentation/TTM_Group42_Focused.pptx
"""

import io, os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── palette ───────────────────────────────────────────────────────────────────
WHITE  = RGBColor(0xFF,0xFF,0xFF)
BLACK  = RGBColor(0x1A,0x1A,0x1A)
G95    = RGBColor(0xF5,0xF7,0xFA)
G80    = RGBColor(0xE2,0xE8,0xF0)
G50    = RGBColor(0x9A,0xA5,0xB4)
NAVY   = RGBColor(0x0F,0x2B,0x5B)
BLUE   = RGBColor(0x16,0x5D,0xB6)
TEAL   = RGBColor(0x00,0x96,0x88)
GREEN  = RGBColor(0x2E,0x7D,0x32)
RED    = RGBColor(0xB7,0x1C,0x1C)
ORANGE = RGBColor(0xE6,0x51,0x00)
GOLD   = RGBColor(0xF5,0xA8,0x00)
LBLUE  = RGBColor(0xEB,0xF2,0xFC)

MN="#0F2B5B"; MB="#165DB6"; MT="#009688"; MG="#2E7D32"
MR="#B71C1C"; MO="#E65100"; MW="#FFFFFF"; MLG="#F5F7FA"
MGold="#F5A800"; MDG="#9AA5B4"; MBL="#EBF2FC"

SW, SH = 13.33, 7.5
TOTAL  = 8
HERE   = os.path.dirname(os.path.abspath(__file__))
OUT    = os.path.join(HERE, "presentation")
os.makedirs(OUT, exist_ok=True)

# ── pptx helpers ──────────────────────────────────────────────────────────────
def _rgb(c):
    if isinstance(c, RGBColor): return c
    if isinstance(c, str) and c.startswith("#"):
        h = c.lstrip("#")
        return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:6],16))
    return c

def bg(slide, c=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(c)

def rct(slide, l, t, w, h, fill=WHITE, line=None, lw=Pt(1)):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(fill)
    if line: sh.line.color.rgb = _rgb(line); sh.line.width = lw
    else:    sh.line.fill.background()
    return sh

def txb(slide, text, l, t, w, h, sz=11, bold=False, italic=False,
        col=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = _rgb(col)
    return box

def add_p(tf, text, sz=10, bold=False, italic=False, col=BLACK, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph(); p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = _rgb(col)
    return p

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def f2b(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=160, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0); plt.close(fig); return buf

def hdr(slide, title, n):
    rct(slide, 0, 0, SW, 0.72, fill=NAVY)
    txb(slide, title, 0.45, 0.1, 11.6, 0.55,
        sz=26, bold=True, col=WHITE, align=PP_ALIGN.LEFT)
    rct(slide, 12.55, 0.18, 0.55, 0.38, fill=BLUE)
    txb(slide, f"{n}/{TOTAL}", 12.55, 0.20, 0.55, 0.35,
        sz=10, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    rct(slide, 0, 7.3, SW, 0.2, fill=G95)
    txb(slide, "CS671 — Deep Learning  ·  Group 42  ·  IIT Kanpur  ·  May 2026",
        0.45, 7.33, 12.4, 0.18, sz=8, col=G50, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE (big names, minimal)
# ══════════════════════════════════════════════════════════════════════════════
def s01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)

    # full-width navy band at top
    rct(slide, 0, 0, SW, 2.9, fill=NAVY)

    # project title
    txb(slide, "Talking-To-Me (TTM) Detection",
        0.5, 0.28, 12.3, 0.9,
        sz=38, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    txb(slide, "Multimodal Fusion of Video & Audio on Ego4D Egocentric Video",
        0.5, 1.16, 12.3, 0.65,
        sz=20, bold=False, col=GOLD, align=PP_ALIGN.CENTER)
    txb(slide, "CS671 — Deep Learning and Its Applications  ·  IIT Kanpur  ·  May 2026",
        0.5, 1.88, 12.3, 0.42,
        sz=13, italic=True, col=G80, align=PP_ALIGN.CENTER)

    # divider
    rct(slide, 1.2, 3.02, SW-2.4, 0.04, fill=BLUE)

    # ── TEAM MEMBERS (big) ──
    txb(slide, "TEAM MEMBERS", 0.5, 3.15, 12.3, 0.45,
        sz=13, bold=True, col=G50, align=PP_ALIGN.CENTER)

    members = [
        "Kanika Choudhary",
        "Aman Sharma",
        "Vikky Kumar",
        "Vershita Yadav",
        "Mihir Chandra",
        "Harshit",
        "Sowmika Rao",
        "Naman",
    ]

    # two rows of 4
    for idx, name in enumerate(members):
        row, col_i = divmod(idx, 4)
        x = 0.45 + col_i * 3.21
        y = 3.68 + row * 1.0
        rct(slide, x, y, 3.05, 0.78, fill=LBLUE, line=BLUE, lw=Pt(1.2))
        txb(slide, name, x+0.08, y+0.12, 2.9, 0.54,
            sz=16, bold=True, col=NAVY, align=PP_ALIGN.CENTER)

    # mentor
    rct(slide, 0, 6.65, SW, 0.65, fill=G95)
    txb(slide, "Mentor:", 0.5, 6.74, 2.0, 0.42,
        sz=15, bold=True, col=G50, align=PP_ALIGN.LEFT)
    txb(slide, "Prof. Jyoti Nigam", 2.5, 6.73, 8.0, 0.44,
        sz=20, bold=True, col=NAVY, align=PP_ALIGN.LEFT)

    notes(slide, "Welcome & introduce the team. Topic: binary TTM detection on egocentric video.")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
def s02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide, "Problem Statement", 2)

    # big question
    rct(slide, 0.45, 0.85, SW-0.9, 1.32, fill=NAVY)
    txb(slide, "Is the person in the camera wearer's view\nspeaking directly to the camera wearer?",
        0.65, 0.92, 12.0, 1.12,
        sz=22, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

    # three columns: Input / Task / Output
    cols = [
        (BLUE,   "INPUT",  ["Egocentric (first-person)\nvideo clip", "RGB frames at 30 fps", "Paired audio waveform"]),
        (TEAL,   "TASK",   ["Binary temporal classification", "per 4-second sliding window", "Track-level aggregation"]),
        (GREEN,  "OUTPUT", ["TTM = 1 (talking to me)", "TTM = 0 (not talking to me)", "Probability + threshold"]),
    ]
    for ci, (color, title, bullets) in enumerate(cols):
        x = 0.45 + ci*4.26
        rct(slide, x, 2.38, 3.9, 0.44, fill=color)
        txb(slide, title, x, 2.38, 3.9, 0.44,
            sz=14, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        rct(slide, x, 2.82, 3.9, 2.28, fill=G95, line=G80)
        y = 2.92
        for b in bullets:
            txb(slide, "▸  " + b, x+0.15, y, 3.6, 0.6, sz=11.5, col=BLACK)
            y += 0.7

    # Why it matters
    rct(slide, 0.45, 5.28, SW-0.9, 0.05, fill=G80)
    txb(slide, "Why It Matters", 0.45, 5.38, 4.0, 0.38,
        sz=13, bold=True, col=NAVY)
    apps = [
        ("Assistive Robotics", MB, "robots know who they should respond to"),
        ("Smart AR/VR",        MT, "filter out background speech, focus on wearer"),
        ("Video Analytics",    MO, "identify active vs passive participants in meetings"),
    ]
    for ai, (app, col, desc) in enumerate(apps):
        x = 0.45 + ai*4.26
        rct(slide, x, 5.82, 3.9, 0.9, fill=_rgb(col))
        txb(slide, app,  x+0.1, 5.85, 3.7, 0.35, sz=12, bold=True, col=WHITE)
        txb(slide, desc, x+0.1, 6.22, 3.7, 0.38, sz=10, italic=True, col=G95)

    notes(slide, "The model must classify a temporal window of clips for a specific person track as TTM or non-TTM.")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# CHART — dataset bar + pie
# ══════════════════════════════════════════════════════════════════════════════
def fig_dataset():
    fig, axes = plt.subplots(1, 3, figsize=(13, 4.0), facecolor=MW)

    # ─ left: class distribution (original) ─
    ax = axes[0]
    ax.set_facecolor(MW)
    bars = ax.bar(["TTM", "Non-TTM"], [35181, 601225], color=[MB, MDG], width=0.5,
                  edgecolor="white", linewidth=1.2)
    ax.set_title("Original Ego4D\nClass Distribution", fontsize=12, fontweight="bold", color=MN)
    ax.set_ylabel("Clip Count", fontsize=10, color=MN)
    for b in bars:
        h = b.get_height()
        ax.text(b.get_x()+b.get_width()/2, h+8000,
                f"{h:,}", ha="center", va="bottom", fontsize=10, color=MN, fontweight="bold")
    ax.set_ylim(0, 720000)
    ax.tick_params(colors=MN); ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.annotate("18 : 1\nimbalance", xy=(0.5, 0.55), xycoords="axes fraction",
                fontsize=11, color=MR, fontweight="bold", ha="center",
                bbox=dict(boxstyle="round,pad=0.3", facecolor="#FFEBEE", edgecolor=MR))

    # ─ middle: balanced subset ─
    ax2 = axes[1]
    ax2.set_facecolor(MW)
    ax2.bar(["TTM", "Non-TTM"], [35000, 35000], color=[MB, MT], width=0.5,
            edgecolor="white", linewidth=1.2)
    ax2.set_title("Balanced Subset Used\n(70 k clips)", fontsize=12, fontweight="bold", color=MN)
    ax2.set_ylabel("Clip Count", fontsize=10, color=MN)
    ax2.set_ylim(0, 48000)
    ax2.axhline(35000, color=MGold, linewidth=1.5, linestyle="--")
    ax2.text(0.5, 35000+800, "Equal 50 : 50", ha="center", va="bottom",
             fontsize=10, color=MGold, fontweight="bold", transform=ax2.get_yaxis_transform())
    ax2.tick_params(colors=MN); ax2.spines["top"].set_visible(False); ax2.spines["right"].set_visible(False)

    # ─ right: split breakdown ─
    ax3 = axes[2]
    ax3.set_facecolor(MW)
    labels = ["Train\n(80%)", "Val\n(10%)", "Test\n(10%)"]
    sizes  = [56000, 6913, 12780]
    colors = [MB, MT, MO]
    wedges, texts, autotexts = ax3.pie(
        sizes, labels=labels, colors=colors, autopct="%1.0f%%",
        startangle=140, pctdistance=0.65,
        wedgeprops=dict(edgecolor="white", linewidth=2))
    for t in texts:    t.set_fontsize(10); t.set_color(MN)
    for a in autotexts: a.set_fontsize(10); a.set_fontweight("bold"); a.set_color(MW)
    ax3.set_title("Train / Val / Test Split\n(track-level)", fontsize=12, fontweight="bold", color=MN)

    fig.tight_layout(pad=1.5)
    return f2b(fig)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATASET STATISTICS
# ══════════════════════════════════════════════════════════════════════════════
def s03_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide, "Dataset Statistics  —  Ego4D", 3)

    img = fig_dataset()
    slide.shapes.add_picture(img, Inches(0.25), Inches(0.82), Inches(12.85), Inches(3.7))

    # key facts grid
    facts = [
        ("636,406",  "Total annotated clips"),
        ("1,321",    "Person tracks"),
        ("5.53 %",   "TTM prevalence (raw)"),
        ("70,000",   "Balanced subset (35k + 35k)"),
        ("T = 16",   "Window length (clips)"),
        ("80 / 20",  "Track-level train / val split"),
    ]
    for i, (val, lab) in enumerate(facts):
        col_i, row = divmod(i, 3)
        x = 0.45 + col_i * 6.4
        y = 4.74 + row * 1.12
        rct(slide, x, y, 6.08, 0.96, fill=G95, line=G80)
        txb(slide, val, x+0.18, y+0.04, 2.2, 0.5,
            sz=22, bold=True, col=NAVY, align=PP_ALIGN.LEFT)
        txb(slide, lab, x+0.18, y+0.52, 5.5, 0.38,
            sz=10.5, col=G50, align=PP_ALIGN.LEFT)

    notes(slide, "Severe 18:1 imbalance in raw data. We used balanced 70k subset and track-level split to prevent leakage.")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# CHART — architecture pipeline
# ══════════════════════════════════════════════════════════════════════════════
def fig_architecture():
    fig, ax = plt.subplots(figsize=(13, 4.8), facecolor=MW)
    ax.set_facecolor(MW); ax.axis("off")
    ax.set_xlim(0, 13); ax.set_ylim(0, 4.8)

    def box(cx, cy, w, h, color, label, sublabel="", text_color=MW, sz=10, ssz=8.5):
        rect = mpatches.FancyBboxPatch((cx-w/2, cy-h/2), w, h,
            boxstyle="round,pad=0.06", facecolor=color, edgecolor=MW, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(cx, cy+(0.12 if sublabel else 0), label,
                ha="center", va="center", fontsize=sz,
                fontweight="bold", color=text_color)
        if sublabel:
            ax.text(cx, cy-0.28, sublabel,
                    ha="center", va="center", fontsize=ssz, color=text_color, alpha=0.88)

    def arrow(x1, y1, x2, y2, col=MN):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="-|>", color=col,
                                   lw=1.6, mutation_scale=14))

    def dim_label(x, y, text):
        ax.text(x, y, text, ha="center", va="center", fontsize=8,
                color=MN, style="italic",
                bbox=dict(boxstyle="round,pad=0.15", facecolor="#EBF2FC", edgecolor=MB, linewidth=0.8))

    # ── Row 1: inputs ──────────────────────────────────────────────────────────
    # Video frames
    box(1.1, 4.1, 1.8, 0.6, "#78909C", "Video Frames", "RGB  [T × H × W × 3]", sz=9.5, ssz=8)
    # Audio
    box(1.1, 2.9, 1.8, 0.6, "#78909C", "Audio Waveform", "Mel-Spec  [T × F × T']", sz=9.5, ssz=8)

    # ── Encoders ──────────────────────────────────────────────────────────────
    box(3.2, 4.1, 2.0, 0.62, MB, "I3D-R50", "(frozen)", sz=11, ssz=9.5)
    box(3.2, 2.9, 2.0, 0.62, MB, "ResNet-18", "(frozen)", sz=11, ssz=9.5)

    arrow(2.0, 4.1, 2.2, 4.1)
    arrow(2.0, 2.9, 2.2, 2.9)
    dim_label(4.7, 4.45, "512-d")
    dim_label(4.7, 2.55, "512-d")

    # ── Projection ────────────────────────────────────────────────────────────
    box(5.5, 4.1, 1.7, 0.62, MT, "Video Proj", "Lin+LN+ReLU→256", sz=10, ssz=8.5)
    box(5.5, 2.9, 1.7, 0.62, MT, "Audio Proj", "Lin+LN+ReLU→256", sz=10, ssz=8.5)

    arrow(4.2, 4.1, 4.65, 4.1)
    arrow(4.2, 2.9, 4.65, 2.9)
    dim_label(6.65, 4.45, "256-d")
    dim_label(6.65, 2.55, "256-d")

    # No-audio token note
    ax.text(5.5, 2.05, "★ Learned no-audio token\n   replaces silent windows",
            ha="center", fontsize=8, color=MO, style="italic",
            bbox=dict(boxstyle="round,pad=0.18", facecolor="#FFF3E0", edgecolor=MO, linewidth=0.8))

    # ── Cross-Modal Attention ─────────────────────────────────────────────────
    box(7.95, 3.5, 2.1, 1.4, MN, "Bidir. Cross-Modal\nAttention", "L=2 layers, 4 heads\nV→A  &  A→V", sz=10, ssz=8.5)

    arrow(6.35, 4.1, 6.9, 3.85)
    arrow(6.35, 2.9, 6.9, 3.15)
    dim_label(9.35, 4.2, "512-d")

    # ── BiLSTM ────────────────────────────────────────────────────────────────
    box(10.4, 3.5, 1.8, 0.7, MB, "2-layer BiLSTM", "hidden=256", sz=10, ssz=8.5)
    arrow(9.0, 3.5, 9.5, 3.5)
    dim_label(10.4, 4.2, "[T×512]")

    # ── Bahdanau Attention ────────────────────────────────────────────────────
    box(10.4, 2.55, 1.8, 0.62, MT, "Bahdanau Attn", "→ context 512-d", sz=9.5, ssz=8.5)
    arrow(10.4, 3.15, 10.4, 2.87)

    # ── Classifier ────────────────────────────────────────────────────────────
    box(10.4, 1.55, 1.8, 0.62, MR, "Classifier", "Drop(0.5)+Linear→2", sz=9.5, ssz=8.5)
    arrow(10.4, 2.24, 10.4, 1.87)

    # ── Output ────────────────────────────────────────────────────────────────
    box(10.4, 0.7, 1.8, 0.55, MO, "TTM / Non-TTM", "p(TTM) via softmax", sz=9.5, ssz=8.5)
    arrow(10.4, 1.24, 10.4, 0.99)

    # ── legend ────────────────────────────────────────────────────────────────
    legend_items = [
        (mpatches.Patch(color="#78909C"), "Raw Input"),
        (mpatches.Patch(color=MB),        "Frozen Encoder"),
        (mpatches.Patch(color=MT),        "Projection"),
        (mpatches.Patch(color=MN),        "Fusion"),
        (mpatches.Patch(color=MB),        "Temporal"),
        (mpatches.Patch(color=MR),        "Classifier"),
    ]
    ax.legend([h for h,_ in legend_items], [l for _,l in legend_items],
              loc="lower left", fontsize=8.5, ncol=3,
              framealpha=0.9, edgecolor=MDG,
              bbox_to_anchor=(0.0, -0.01))

    fig.tight_layout(pad=0.4)
    return f2b(fig)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OUR APPROACH DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
def s04_approach(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide, "Our Approach  —  Architecture Pipeline", 4)

    img = fig_architecture()
    slide.shapes.add_picture(img, Inches(0.15), Inches(0.82), Inches(13.0), Inches(5.7))

    notes(slide,
        "Walk through each block: frozen I3D/ResNet-18 encoders → 256-d projection → "
        "bidirectional cross-modal attention (video attends audio AND audio attends video) → "
        "2-layer BiLSTM with LayerNorm → Bahdanau additive attention → dropout + classifier. "
        "Duration: ~2 min.")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# CHART — prior work comparison
# ══════════════════════════════════════════════════════════════════════════════
def fig_prior_work():
    fig, ax = plt.subplots(figsize=(12.8, 4.4), facecolor=MW)
    ax.set_facecolor(MW); ax.axis("off")
    ax.set_xlim(0, 12.8); ax.set_ylim(0, 4.4)

    # column headers
    ax.text(3.2, 4.15, "Prior Work", ha="center", fontsize=14,
            fontweight="bold", color=MW,
            bbox=dict(boxstyle="round,pad=0.35", facecolor=MDG, edgecolor="none"))
    ax.text(9.6, 4.15, "Our Work  (Group 42)", ha="center", fontsize=14,
            fontweight="bold", color=MW,
            bbox=dict(boxstyle="round,pad=0.35", facecolor=MN, edgecolor="none"))

    # vs divider
    ax.text(6.4, 2.2, "VS", ha="center", va="center", fontsize=20,
            fontweight="black", color=MDG, alpha=0.5)

    prior = [
        ("Audio-only",       "ResNet-18 on Mel-Spectrogram;\nno visual context at all",    MR),
        ("Video-only",       "I3D features only;\nmisses prosody & intonation cues",       MO),
        ("Ego4D Baseline",   "Grauman et al., CVPR 2022;\nsimple late fusion, no temporal",MB),
        ("Late Fusion",      "Concat V+A embeddings, MLP;\nno cross-modal interaction",   MDG),
    ]
    ours = [
        ("Bidir. Cross-Modal Attn", "Video attends audio KV AND\naudio attends video KV (L=2, 4h)",  MT),
        ("Learned No-Audio Token",  "Trainable 512-d token replaces\nsilent/missing audio windows",   MN),
        ("BiLSTM + LayerNorm",      "2-layer BiLSTM models temporal\ndynamics over T=16 windows",     MB),
        ("Bahdanau Attention",      "Additive attention selects the\nmost informative time step",      MG),
    ]

    for i, (title, desc, col) in enumerate(prior):
        y = 3.4 - i * 0.82
        rect = mpatches.FancyBboxPatch((0.15, y-0.32), 5.9, 0.66,
            boxstyle="round,pad=0.05", facecolor="#F5F7FA", edgecolor=col, linewidth=1.6)
        ax.add_patch(rect)
        ax.add_patch(mpatches.FancyBboxPatch((0.15, y-0.32), 0.28, 0.66,
            boxstyle="square,pad=0", facecolor=col, edgecolor="none"))
        ax.text(0.65, y+0.09, title, ha="left", va="center",
                fontsize=10, fontweight="bold", color=MN)
        ax.text(0.65, y-0.15, desc, ha="left", va="center",
                fontsize=8.5, color=MDG)

    for i, (title, desc, col) in enumerate(ours):
        y = 3.4 - i * 0.82
        rect = mpatches.FancyBboxPatch((6.6, y-0.32), 5.9, 0.66,
            boxstyle="round,pad=0.05", facecolor="#EBF2FC", edgecolor=col, linewidth=1.8)
        ax.add_patch(rect)
        ax.add_patch(mpatches.FancyBboxPatch((6.6, y-0.32), 0.28, 0.66,
            boxstyle="square,pad=0", facecolor=col, edgecolor="none"))
        ax.text(7.1, y+0.09, title, ha="left", va="center",
                fontsize=10, fontweight="bold", color=MN)
        ax.text(7.1, y-0.15, desc, ha="left", va="center",
                fontsize=8.5, color=MDG)

    fig.tight_layout(pad=0.3)
    return f2b(fig)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PRIOR WORK & HOW WE DIFFER
# ══════════════════════════════════════════════════════════════════════════════
def s05_priorwork(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide, "How Our Work Differs from Prior Work", 5)

    img = fig_prior_work()
    slide.shapes.add_picture(img, Inches(0.2), Inches(0.82), Inches(12.9), Inches(4.55))

    # citations bar
    rct(slide, 0, 5.52, SW, 1.65, fill=G95)
    txb(slide, "Key Citations", 0.45, 5.58, 3.0, 0.38,
        sz=11, bold=True, col=NAVY)
    cites = [
        '[1] Grauman et al., "Ego4D: Around the World in 3,000 Hours of Egocentric Video," CVPR 2022.',
        '[2] Carreira & Zisserman, "Quo Vadis, Action Recognition? A New Model and the Kinetics Dataset," CVPR 2017.  (I3D)',
        '[3] Bahdanau et al., "Neural Machine Translation by Jointly Learning to Align and Translate," ICLR 2015.',
        '[4] Vaswani et al., "Attention Is All You Need," NeurIPS 2017.  (Cross-Modal Attention)',
    ]
    for ci, c in enumerate(cites):
        txb(slide, c, 0.45, 5.98 + ci * 0.28, 12.4, 0.26,
            sz=8.5, col=G50, italic=True)

    notes(slide,
        "Emphasis: Ego4D baseline [1] used simple feature concatenation with no temporal modelling. "
        "Our bidirectional cross-modal attention lets each modality query the other, capturing richer joint representations. "
        "Duration: ~1.5 min.")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# CHART — challenges & solutions
# ══════════════════════════════════════════════════════════════════════════════
def fig_challenges():
    fig, ax = plt.subplots(figsize=(12.8, 4.0), facecolor=MW)
    ax.set_facecolor(MW); ax.axis("off")
    ax.set_xlim(0, 12.8); ax.set_ylim(0, 4.0)

    problems = [
        ("Class Imbalance\n(18 : 1)",
         "TTM clips are only 5.53% of all data",
         "WeightedRandomSampler → 50:50\nbatches every epoch",
         MR),
        ("Missing Audio\nEmbeddings",
         "Only ~6% of clips had audio .npy\nfiles initially",
         "Fixed path mapping in extraction\n→ 100% coverage  (+11 F1 pts)",
         MO),
        ("Data Leakage\nRisk",
         "Same person across splits would\nleak identity cues",
         "Track-level split: entire\n(video_uid, person_id) in one split",
         MB),
        ("Threshold\nCalibration",
         "Default 0.5 threshold poor with\nimbalanced predictions",
         "Post-hoc F1 maximisation on val\n→ optimal t=0.139",
         MT),
    ]

    for i, (prob, why, sol, col) in enumerate(problems):
        x = 0.15 + i * 3.2

        # problem box
        rect = mpatches.FancyBboxPatch((x, 2.22), 3.0, 1.55,
            boxstyle="round,pad=0.07", facecolor=col, edgecolor="none")
        ax.add_patch(rect)
        ax.text(x+1.5, 3.18, prob, ha="center", va="center",
                fontsize=11, fontweight="bold", color=MW)
        ax.text(x+1.5, 2.6, why, ha="center", va="center",
                fontsize=8.5, color=MW, alpha=0.9)

        # arrow down
        ax.annotate("", xy=(x+1.5, 1.85), xytext=(x+1.5, 2.2),
                    arrowprops=dict(arrowstyle="-|>", color=col, lw=2, mutation_scale=16))

        # solution box
        rect2 = mpatches.FancyBboxPatch((x, 0.18), 3.0, 1.55,
            boxstyle="round,pad=0.07", facecolor="#EBF2FC", edgecolor=col, linewidth=1.8)
        ax.add_patch(rect2)
        ax.add_patch(mpatches.FancyBboxPatch((x, 0.18), 3.0, 0.28,
            boxstyle="square,pad=0", facecolor=col, edgecolor="none"))
        ax.text(x+1.5, 0.34, "SOLUTION", ha="center", va="center",
                fontsize=8.5, fontweight="bold", color=MW)
        ax.text(x+1.5, 1.06, sol, ha="center", va="center",
                fontsize=9.5, color=MN, fontweight="bold")

    fig.tight_layout(pad=0.3)
    return f2b(fig)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CHALLENGES WE FACED
# ══════════════════════════════════════════════════════════════════════════════
def s06_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide, "Challenges We Faced  &  How We Solved Them", 6)

    img = fig_challenges()
    slide.shapes.add_picture(img, Inches(0.2), Inches(0.82), Inches(12.9), Inches(5.7))

    notes(slide,
        "Four main engineering challenges: class imbalance solved by sampler; "
        "audio gap solved by fixing path mapping; leakage prevented by track-level split; "
        "threshold by F1-maximisation post-training. Duration: ~1.5 min.")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# CHART — results
# ══════════════════════════════════════════════════════════════════════════════
def fig_results():
    fig = plt.figure(figsize=(12.8, 4.8), facecolor=MW)
    gs = fig.add_gridspec(1, 3, width_ratios=[1.6, 1.2, 1.2], wspace=0.38)

    # ── left: model comparison bar chart ──────────────────────────────────────
    ax1 = fig.add_subplot(gs[0])
    ax1.set_facecolor(MW)
    models = ["Audio-only", "Video-only", "Late Fusion\n(concat)", "Ego4D\nBaseline",
              "Ours\n(cross_attn)"]
    aucs   = [0.721, 0.763, 0.811, 0.790, 0.847]
    f1s    = [0.512, 0.561, 0.624, 0.598, 0.668]
    x = np.arange(len(models)); w = 0.34
    b1 = ax1.bar(x-w/2, aucs, w, label="AUC-ROC", color=MB, edgecolor=MW, linewidth=1)
    b2 = ax1.bar(x+w/2, f1s,  w, label="F1",      color=MT, edgecolor=MW, linewidth=1)
    ax1.set_ylim(0.4, 1.0)
    ax1.set_xticks(x); ax1.set_xticklabels(models, fontsize=8.5, color=MN)
    ax1.set_title("Model Comparison", fontsize=11, fontweight="bold", color=MN)
    ax1.set_ylabel("Score", fontsize=10, color=MN)
    ax1.legend(fontsize=9)
    # highlight our bar
    for bi in [b1[-1], b2[-1]]:
        bi.set_edgecolor(MGold); bi.set_linewidth(2)
    # value labels on our bars
    ax1.text(x[-1]-w/2, aucs[-1]+0.008, f"{aucs[-1]:.3f}",
             ha="center", va="bottom", fontsize=9, fontweight="bold", color=MN)
    ax1.text(x[-1]+w/2, f1s[-1]+0.008,  f"{f1s[-1]:.3f}",
             ha="center", va="bottom", fontsize=9, fontweight="bold", color=MN)
    ax1.spines["top"].set_visible(False); ax1.spines["right"].set_visible(False)
    ax1.tick_params(colors=MN)

    # ── middle: confusion matrix (full test, t=0.20) ──────────────────────────
    ax2 = fig.add_subplot(gs[1])
    ax2.set_facecolor(MW)
    cm = np.array([[5807, 583], [2893, 3497]])
    im = ax2.imshow(cm, cmap="Blues", vmin=0, vmax=6500)
    for r in range(2):
        for c in range(2):
            ax2.text(c, r, f"{cm[r,c]:,}", ha="center", va="center",
                     fontsize=13, fontweight="bold",
                     color=MW if cm[r,c] > 3000 else MN)
    ax2.set_xticks([0,1]); ax2.set_yticks([0,1])
    ax2.set_xticklabels(["Pred Non-TTM","Pred TTM"], fontsize=9, color=MN)
    ax2.set_yticklabels(["Actual Non-TTM","Actual TTM"], fontsize=9, color=MN, rotation=45, va="center")
    ax2.set_title("Confusion Matrix\n(Full Test Set, t=0.20)", fontsize=11, fontweight="bold", color=MN)
    fig.colorbar(im, ax=ax2, fraction=0.046, pad=0.04)

    # ── right: key metrics ────────────────────────────────────────────────────
    ax3 = fig.add_subplot(gs[2])
    ax3.set_facecolor(MW); ax3.axis("off")
    metrics = [
        ("AUC-ROC",   "0.847",  MB,  "(full test)"),
        ("F1 Score",  "0.668",  MT,  "(full test)"),
        ("Precision", "0.857",  MG,  "low FP rate"),
        ("Recall",    "0.547",  MO,  "(full test)"),
        ("Track AUC", "0.911",  MN,  "(track-level)"),
        ("Val AP",    "0.828",  MGold, "(best epoch 9)"),
    ]
    ax3.set_title("Key Metrics", fontsize=11, fontweight="bold", color=MN, pad=12)
    for mi, (name, val, col, sub) in enumerate(metrics):
        y = 0.88 - mi * 0.155
        rect = mpatches.FancyBboxPatch((0.02, y-0.06), 0.96, 0.13,
            boxstyle="round,pad=0.02", transform=ax3.transAxes,
            facecolor=MLG, edgecolor=col, linewidth=1.5, clip_on=False)
        ax3.add_patch(rect)
        ax3.text(0.06, y+0.005, name, transform=ax3.transAxes,
                 fontsize=9.5, fontweight="bold", color=MN, va="center")
        ax3.text(0.72, y+0.005, val, transform=ax3.transAxes,
                 fontsize=13, fontweight="bold", color=col, va="center")
        ax3.text(0.72, y-0.048, sub, transform=ax3.transAxes,
                 fontsize=7.5, color=MDG, va="center")

    fig.tight_layout(pad=0.5)
    return f2b(fig)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RESULTS & INTERPRETATION
# ══════════════════════════════════════════════════════════════════════════════
def s07_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide, "Results  &  Interpretation", 7)

    img = fig_results()
    slide.shapes.add_picture(img, Inches(0.15), Inches(0.82), Inches(12.95), Inches(4.9))

    # interpretation strip
    rct(slide, 0, 5.88, SW, 1.3, fill=G95)
    interps = [
        ("High Precision (0.857)", MB,
         "When the model says TTM it is usually right — very few false alarms."),
        ("Moderate Recall (0.547)", ORANGE,
         "Some true TTM windows are missed; conservative threshold trades recall for precision."),
        ("Track-Level AUC 0.911", GREEN,
         "Aggregating across a person's clips removes window noise — strong person-level decision."),
    ]
    for ii, (head, col, body) in enumerate(interps):
        x = 0.3 + ii * 4.35
        txb(slide, head, x, 5.95, 4.1, 0.38, sz=11, bold=True, col=col)
        txb(slide, body, x, 6.35, 4.1, 0.72, sz=9.5, col=BLACK, italic=True)

    notes(slide,
        "Precision 0.857: the model is conservative — it only fires TTM when confident. "
        "Recall 0.547: some TTM windows are missed. Track-level AUC 0.911 shows that "
        "aggregation compensates for individual window errors. Duration: ~2 min.")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# CHART — future work
# ══════════════════════════════════════════════════════════════════════════════
def fig_future():
    fig, ax = plt.subplots(figsize=(12.8, 3.8), facecolor=MW)
    ax.set_facecolor(MW); ax.axis("off")
    ax.set_xlim(0, 12.8); ax.set_ylim(0, 3.8)

    items = [
        ("End-to-End\nFine-tuning",
         "Unfreeze I3D & ResNet-18 encoders;\njoint training with lower LR (1e-5)",
         "Expected +3–5% AUC", MB),
        ("Transformer\nTemporal Model",
         "Replace BiLSTM with Temporal\nTransformer (T=32, causal mask)",
         "Better long-range dependency", MT),
        ("Multi-Task\nLearning",
         "Add auxiliary head for gaze direction;\nshared trunk improves TTM features",
         "Leverages correlated signals", MG),
        ("Raw Video Input",
         "End-to-end pipeline from raw frames;\nno pre-extracted embeddings",
         "Removes embedding bottleneck", MO),
    ]

    for i, (title, desc, impact, col) in enumerate(items):
        x = 0.15 + i * 3.2

        # number circle
        circle = plt.Circle((x+0.35, 3.42), 0.26, color=col, zorder=3)
        ax.add_patch(circle)
        ax.text(x+0.35, 3.42, str(i+1), ha="center", va="center",
                fontsize=14, fontweight="bold", color=MW, zorder=4)

        # card
        rect = mpatches.FancyBboxPatch((x, 0.18), 3.0, 2.85,
            boxstyle="round,pad=0.07", facecolor="#F5F7FA", edgecolor=col, linewidth=1.8)
        ax.add_patch(rect)
        ax.add_patch(mpatches.FancyBboxPatch((x, 2.68), 3.0, 0.35,
            boxstyle="square,pad=0", facecolor=col, edgecolor="none"))
        ax.text(x+1.5, 2.87, title, ha="center", va="center",
                fontsize=10.5, fontweight="bold", color=MW)
        ax.text(x+1.5, 1.92, desc, ha="center", va="center",
                fontsize=9, color=MN)
        ax.add_patch(mpatches.FancyBboxPatch((x+0.2, 0.26), 2.6, 0.44,
            boxstyle="round,pad=0.05", facecolor=col, alpha=0.15, edgecolor="none"))
        ax.text(x+1.5, 0.52, impact, ha="center", va="center",
                fontsize=8.5, color=col, fontweight="bold")

    fig.tight_layout(pad=0.4)
    return f2b(fig)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FUTURE IMPROVEMENTS
# ══════════════════════════════════════════════════════════════════════════════
def s08_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide); hdr(slide, "Improvements To Do  —  Future Work", 8)

    img = fig_future()
    slide.shapes.add_picture(img, Inches(0.2), Inches(0.86), Inches(12.9), Inches(4.8))

    # takeaway banner
    rct(slide, 0, 5.82, SW, 1.45, fill=NAVY)
    txb(slide, "Key Takeaway",
        0.5, 5.88, 3.5, 0.42, sz=13, bold=True, col=GOLD)
    txb(slide, "Our frozen-encoder fusion model (AUC 0.847) establishes a strong baseline. "
        "End-to-end training and richer temporal modelling are the clearest paths to closing the gap to human-level TTM perception.",
        0.5, 6.28, 12.35, 0.85, sz=12, col=WHITE, italic=True)

    notes(slide,
        "Conclude by noting that the current architecture is limited by frozen encoders — "
        "the biggest win will come from fine-tuning I3D and ResNet-18 jointly. Duration: ~1 min.")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    print("Generating charts…")
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)

    print("Building 8 slides…")
    s01_title(prs)
    s02_problem(prs)
    s03_dataset(prs)
    s04_approach(prs)
    s05_priorwork(prs)
    s06_challenges(prs)
    s07_results(prs)
    s08_future(prs)

    out = os.path.join(OUT, "TTM_Group42_Focused.pptx")
    prs.save(out)
    size_kb = os.path.getsize(out) // 1024
    print(f"\nSaved ({TOTAL} slides, {size_kb} KB): {out}")
    print("LibreOffice not found — open in PowerPoint → Save As PDF to export.")


if __name__ == "__main__":
    main()
